#!/usr/bin/env python3
"""
Generate a 10-slide NVIDIA-branded overview deck for introducing
the Agentic AI Science Playbook to an internal NVIDIA team.

Reuses helpers and branding from the main deck generator.

Usage:
    python generate_overview_deck.py

Output:
    Playbook_Team_Overview.pptx
"""

from __future__ import annotations
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ── NVIDIA Brand ────────────────────────────────────────────────────────────

NVIDIA_GREEN  = RGBColor(0x76, 0xB9, 0x00)
NVIDIA_BLACK  = RGBColor(0x1A, 0x1A, 0x1A)
NVIDIA_DARK   = RGBColor(0x24, 0x24, 0x24)
NVIDIA_GRAY   = RGBColor(0x2D, 0x2D, 0x2D)
NVIDIA_LGRAY  = RGBColor(0x8C, 0x8C, 0x8C)
NVIDIA_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
NVIDIA_GREEN2 = RGBColor(0x5A, 0x9E, 0x00)
CLR_EOP       = RGBColor(0x88, 0x44, 0xCC)
CLR_HEALTH    = RGBColor(0xCC, 0x44, 0x44)
CLR_BIO       = RGBColor(0x00, 0x88, 0xCC)
CLR_FINANCE   = RGBColor(0x00, 0x88, 0x55)

FONT = "Calibri"
FONT_CODE = "Consolas"
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NVIDIA_BLACK

def _bar(slide):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(0.08))
    s.fill.solid(); s.fill.fore_color.rgb = NVIDIA_GREEN; s.line.fill.background()
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), SLIDE_H - Inches(0.06), SLIDE_W, Inches(0.06))
    s.fill.solid(); s.fill.fore_color.rgb = NVIDIA_GREEN; s.line.fill.background()

def _label(slide, text):
    t = slide.shapes.add_textbox(Inches(0.6), Inches(0.22), Inches(5), Inches(0.3))
    p = t.text_frame.paragraphs[0]
    p.text = text.upper(); p.font.size = Pt(11); p.font.color.rgb = NVIDIA_GREEN
    p.font.bold = True; p.font.name = FONT

def _title(slide, text, top=Inches(0.55), size=Pt(36)):
    t = slide.shapes.add_textbox(Inches(0.6), top, Inches(12), Inches(0.8))
    t.text_frame.word_wrap = True
    p = t.text_frame.paragraphs[0]
    p.text = text; p.font.size = size; p.font.color.rgb = NVIDIA_WHITE
    p.font.bold = True; p.font.name = FONT

def _sub(slide, text, top=Inches(1.25)):
    t = slide.shapes.add_textbox(Inches(0.6), top, Inches(11), Inches(0.5))
    t.text_frame.word_wrap = True
    p = t.text_frame.paragraphs[0]
    p.text = text; p.font.size = Pt(18); p.font.color.rgb = NVIDIA_LGRAY; p.font.name = FONT

def _box(slide, text, left, top, w, h, fill=NVIDIA_GRAY, color=NVIDIA_WHITE, size=Pt(14), bold=False):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill; s.line.fill.background()
    tf = s.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.1); tf.margin_top = Inches(0.05)
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        p.text = line; p.font.size = size; p.font.color.rgb = color; p.font.bold = bold; p.font.name = FONT
    return s

def _arrow_r(slide, l, t, w=Inches(0.4), h=Inches(0.25)):
    s = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = NVIDIA_GREEN; s.line.fill.background()

def _arrow_d(slide, l, t, w=Inches(0.3), h=Inches(0.35)):
    s = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = NVIDIA_GREEN; s.line.fill.background()

def _table(slide, headers, rows, left=Inches(0.6), top=Inches(2.0), width=Inches(12)):
    rh = Inches(0.42)
    nr, nc = len(rows)+1, len(headers)
    cw = width // nc
    ts = slide.shapes.add_table(nr, nc, left, top, width, rh*nr)
    t = ts.table
    for j, h in enumerate(headers):
        c = t.cell(0, j); c.text = h
        for p in c.text_frame.paragraphs:
            p.font.size = Pt(13); p.font.bold = True; p.font.color.rgb = NVIDIA_BLACK; p.font.name = FONT
        c.fill.solid(); c.fill.fore_color.rgb = NVIDIA_GREEN
    for i, row in enumerate(rows):
        for j, v in enumerate(row):
            c = t.cell(i+1, j); c.text = v
            for p in c.text_frame.paragraphs:
                p.font.size = Pt(12); p.font.color.rgb = NVIDIA_WHITE; p.font.name = FONT
            c.fill.solid(); c.fill.fore_color.rgb = NVIDIA_DARK if i%2==0 else NVIDIA_GRAY
    for j in range(nc):
        t.columns[j].width = int(cw)

def _notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text

def _common(slide, label=None):
    _bg(slide); _bar(slide)
    if label: _label(slide, label)


# ═══════════════════════════════════════════════════════════════════════════
# 10 SLIDES
# ═══════════════════════════════════════════════════════════════════════════

def slide_01(prs):
    """Title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(0.15))
    s.fill.solid(); s.fill.fore_color.rgb = NVIDIA_GREEN; s.line.fill.background()

    t = slide.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(3), Inches(0.5))
    p = t.text_frame.paragraphs[0]
    p.text = "NVIDIA"; p.font.size = Pt(18); p.font.color.rgb = NVIDIA_GREEN; p.font.bold = True; p.font.name = FONT

    _title(slide, "Agentic AI for\nScientific Discovery", top=Inches(2.0), size=Pt(48))

    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(4.0), Inches(3), Inches(0.04))
    ln.fill.solid(); ln.fill.fore_color.rgb = NVIDIA_GREEN; ln.line.fill.background()

    t2 = slide.shapes.add_textbox(Inches(0.6), Inches(4.3), Inches(11), Inches(0.5))
    t2.text_frame.word_wrap = True
    p2 = t2.text_frame.paragraphs[0]
    p2.text = "A Hands-On Playbook — Team Overview"
    p2.font.size = Pt(24); p2.font.color.rgb = NVIDIA_LGRAY; p2.font.name = FONT

    t3 = slide.shapes.add_textbox(Inches(0.6), Inches(5.2), Inches(11), Inches(0.5))
    p3 = t3.text_frame.paragraphs[0]
    p3.text = "20 Labs  |  5 Domains  |  NVIDIA NIM  |  Beginner-Friendly"
    p3.font.size = Pt(16); p3.font.color.rgb = NVIDIA_WHITE; p3.font.name = FONT

    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), SLIDE_H - Inches(0.06), SLIDE_W, Inches(0.06))
    s.fill.solid(); s.fill.fore_color.rgb = NVIDIA_GREEN; s.line.fill.background()

    _notes(slide,
        "Team — I want to walk you through a playbook we've built that teaches scientists how to "
        "build AI agents for their research. 20 hands-on labs, 5 scientific domains, all running on "
        "NVIDIA NIM. This is a 10-minute overview — I'll cover what it is, who it's for, what's in it, "
        "and how we can use it with customers.\n\nLet's dive in.")


def slide_02(prs):
    """The Problem — why agents matter."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _common(slide, "THE PROBLEM")
    _title(slide, "Scientists Spend 80% of Their Time on the Wrong Things")

    pairs = [
        ("Data wrangling", "weeks", "AI agent automates", "minutes"),
        ("Literature review", "2-3 weeks", "RAG agent searches + synthesizes", "30 min"),
        ("Manual tool-switching", "10+ interfaces", "Agent orchestrates all tools", "1 interface"),
        ("Writing protocols", "hours per protocol", "Agent generates + validates", "seconds"),
        ("Quality checking", "subjective review", "LLM-as-Judge auto-scores", "systematic"),
    ]
    for i, (task, before, solution, after) in enumerate(pairs):
        tp = Inches(1.7) + Inches(i * 1.05)
        _box(slide, task, Inches(0.4), tp, Inches(2.5), Inches(0.75), fill=RGBColor(0x8B, 0x00, 0x00), size=Pt(13), bold=True)
        _box(slide, before, Inches(3.0), tp, Inches(1.8), Inches(0.75), fill=NVIDIA_DARK, color=NVIDIA_LGRAY, size=Pt(12))
        _arrow_r(slide, Inches(4.9), tp + Inches(0.22), Inches(0.4), Inches(0.25))
        _box(slide, solution, Inches(5.4), tp, Inches(4.5), Inches(0.75), fill=NVIDIA_GREEN2, size=Pt(12), bold=True)
        _box(slide, after, Inches(10.1), tp, Inches(2.5), Inches(0.75), fill=NVIDIA_GREEN, color=NVIDIA_BLACK, size=Pt(13), bold=True)

    _notes(slide,
        "This is the core pitch. Scientists are spending 80% of their time on repetitive tasks — "
        "data wrangling, literature search, manual tool-switching. AI agents solve each of these "
        "by taking ACTION, not just generating text.\n\nLook at literature review: 2-3 weeks manually "
        "vs. 30 minutes with a RAG agent. That's the kind of ROI we're offering.\n\nThe playbook "
        "teaches researchers how to build these agents themselves.")


def slide_03(prs):
    """What's in the playbook — structure."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _common(slide, "WHAT'S IN THE PLAYBOOK")
    _title(slide, "20 Labs Across 5 Scientific Domains")
    _sub(slide, "Foundation (universal) + Domain tracks (specialized)")

    # Foundation bar
    _box(slide, "FOUNDATION — 7 Labs + Scenarios Intro", Inches(0.4), Inches(1.9), Inches(12.3), Inches(0.5),
         fill=NVIDIA_GREEN, color=NVIDIA_BLACK, size=Pt(15), bold=True)

    labs = ["Scenarios\n(Intro)", "Lab 0\nPrototype", "Lab 1-2\nSchemas", "Lab 3\nMemory",
            "Lab 4\nGraphs", "Lab 5\nEval", "Lab 6\nCo-Scientist"]
    for i, lab in enumerate(labs):
        _box(slide, lab, Inches(0.4) + Inches(i*1.82), Inches(2.55), Inches(1.62), Inches(0.75),
             fill=NVIDIA_GRAY, size=Pt(10))

    _arrow_d(slide, Inches(6.2), Inches(3.4), Inches(0.35), Inches(0.3))

    # Domain boxes
    domains = [
        ("EOP\n3 Labs", CLR_EOP), ("Healthcare\n3 Labs", CLR_HEALTH),
        ("Bioinformatics\n3 Labs", CLR_BIO), ("Finance\n3 Labs", CLR_FINANCE),
    ]
    for i, (name, color) in enumerate(domains):
        _box(slide, name, Inches(0.4) + Inches(i*3.15), Inches(3.9), Inches(2.85), Inches(0.85),
             fill=color, size=Pt(14), bold=True)

    # Roadmap
    _box(slide, "ROADMAP:  Drug Discovery (BioNeMo)  |  Materials (PhysicsNeMo)  |  Climate (Earth-2)  |  Chemistry (cuOpt)  |  Robotics (Isaac Sim)",
         Inches(0.4), Inches(5.1), Inches(12.3), Inches(0.5), fill=NVIDIA_DARK, color=NVIDIA_LGRAY, size=Pt(12))

    # Stats
    _box(slide, "20 notebooks available today  |  Every lab has step-by-step explanations + expected output  |  Works on any laptop",
         Inches(0.4), Inches(5.9), Inches(12.3), Inches(0.45), fill=NVIDIA_GREEN, color=NVIDIA_BLACK, size=Pt(13), bold=True)

    _notes(slide,
        "Here's the structure. Two layers: Foundation is universal — 7 labs that teach agent engineering "
        "patterns applicable to ANY domain. Then 4 domain tracks with 3 labs each.\n\n"
        "EOP for computational science reproducibility. Healthcare for clinical NLP, literature synthesis, "
        "and trial screening. Bioinformatics for genomics. And Finance for quant analysis, portfolio "
        "optimization, and algo strategies.\n\n"
        "On the roadmap: 5 more domains integrating NVIDIA's domain SDKs. Each follows the same 3-lab template.")


def slide_04(prs):
    """Learning roadmap — the journey."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _common(slide, "LEARNING JOURNEY")
    _title(slide, "From Zero to AI Co-Scientist")
    _sub(slide, "Clear path for beginners — no AI expertise required")

    stages = [
        ("START\nHERE", "Scenarios Lab\n6 use cases\n(no coding)", RGBColor(0x44, 0x44, 0xCC)),
        ("LEARN", "Foundation\nLabs 0-6\n(~8 hours)", NVIDIA_GREEN),
        ("SPECIALIZE", "Pick your\ndomain track\n(~3 hours)", RGBColor(0xCC, 0x88, 0x00)),
        ("BUILD", "Apply to YOUR\nresearch\n(ongoing)", NVIDIA_GREEN),
    ]
    for i, (label, desc, color) in enumerate(stages):
        left = Inches(0.4) + Inches(i * 3.2)
        _box(slide, f"{label}\n\n{desc}", left, Inches(1.9), Inches(2.8), Inches(2.8),
             fill=color, size=Pt(13))
        # Bold the title
        shape = slide.shapes[-1]
        shape.text_frame.paragraphs[0].runs[0].font.bold = True
        shape.text_frame.paragraphs[0].runs[0].font.size = Pt(18)
        if i < 3:
            _arrow_r(slide, left + Inches(2.85), Inches(3.05), Inches(0.3), Inches(0.25))

    # Beginner features
    _box(slide, "Built for beginners:  Step-by-step code explanations  |  Expected output for every experiment  |  "
         "LangGraph warm-up  |  Genomics primer  |  NIM vs OpenAI guide",
         Inches(0.4), Inches(5.2), Inches(12.3), Inches(0.55), fill=NVIDIA_GRAY, color=NVIDIA_WHITE, size=Pt(12))

    _box(slide, "Prerequisites: Python (basic) + curiosity. No AI/ML expertise needed.",
         Inches(0.4), Inches(6.0), Inches(12.3), Inches(0.4), fill=NVIDIA_GREEN, color=NVIDIA_BLACK, size=Pt(14), bold=True)

    _notes(slide,
        "The learning path is designed for people who know some Python but are NOT agent experts.\n\n"
        "Start with the Scenarios Lab — zero coding, just see what agents can do. 45 minutes.\n"
        "Then Foundation — 7 labs over about 8 hours, spread across a few days.\n"
        "Then pick a domain and specialize.\n\n"
        "We've built in beginner-friendly features: every code cell has an explanation, every experiment "
        "has collapsible expected output, and tricky topics like LangGraph and genomics have warm-up sections.")


def slide_05(prs):
    """NVIDIA NIM — why it matters."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _common(slide, "NVIDIA NIM")
    _title(slide, "Why NVIDIA NIM for Scientific Agents")
    _sub(slide, "One environment variable to switch from OpenAI — zero code changes")

    _table(slide,
        ["Factor", "OpenAI API", "NVIDIA NIM"],
        [
            ["Data Privacy", "Data sent to OpenAI", "Stays in YOUR infrastructure"],
            ["Cost at Scale", "Per-token pricing", "Per-GPU-hour (3-10x cheaper at 10K+ calls)"],
            ["Latency", "~500ms", "~200ms (GPU-optimized)"],
            ["Models", "GPT-4o (closed)", "Nemotron, Llama 3.3, 100+ open models"],
            ["Compliance", "Cloud only", "On-prem for HIPAA / pharma / finance"],
        ],
        top=Inches(1.7)
    )

    # Code snippet
    code_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.0), Inches(11.5), Inches(1.0))
    code_box.fill.solid(); code_box.fill.fore_color.rgb = RGBColor(0x1E, 0x1E, 0x2E)
    code_box.line.color.rgb = RGBColor(0x3A, 0x3A, 0x4A); code_box.line.width = Pt(1)
    tf = code_box.text_frame; tf.margin_left = Inches(0.2)
    p = tf.paragraphs[0]
    p.text = '# Switch in one line:  export USE_NIM=true && export NIM_API_KEY="nvapi-..."'
    p.font.size = Pt(14); p.font.color.rgb = NVIDIA_GREEN; p.font.name = FONT_CODE; p.font.bold = True

    _box(slide, "All 20 labs work with both OpenAI and NIM. Free API key at build.nvidia.com",
         Inches(0.4), Inches(6.3), Inches(12.3), Inches(0.4), fill=NVIDIA_GREEN, color=NVIDIA_BLACK, size=Pt(14), bold=True)

    _notes(slide,
        "Every lab in the playbook works with both OpenAI and NVIDIA NIM. The switch is one "
        "environment variable.\n\nThe customer conversation starter: 'What are your data privacy requirements?' "
        "If they handle patient data, proprietary molecules, or financial data — NIM keeps it on-prem.\n\n"
        "At scale, NIM is 3-10x cheaper than per-token API pricing. And Nemotron models are open-weight, "
        "so teams can inspect, fine-tune, and audit them.")


def slide_06(prs):
    """Highlight: AI Co-Scientist multi-agent."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _common(slide, "HIGHLIGHT")
    _title(slide, "Lab 6: AI Co-Scientist — Multi-Agent Research")
    _sub(slide, "Four specialized agents collaborate on scientific discovery")

    # Orchestrator
    _box(slide, "ORCHESTRATOR\nCoordinates the research team", Inches(4.2), Inches(1.8), Inches(4.8), Inches(0.8),
         fill=NVIDIA_GREEN, color=NVIDIA_BLACK, size=Pt(14), bold=True)
    _arrow_d(slide, Inches(6.4), Inches(2.65), Inches(0.3), Inches(0.35))

    agents = [
        ("LITERATURE\nAGENT", "Searches papers\nFinds gaps", CLR_HEALTH),
        ("HYPOTHESIS\nAGENT", "Generates novel\nideas from gaps", RGBColor(0x00, 0xAA, 0x88)),
        ("CRITIC\nAGENT", "Scores novelty\nfeasibility, impact", RGBColor(0xCC, 0x88, 0x00)),
        ("EXPERIMENT\nAGENT", "Designs protocols\nto test hypotheses", CLR_BIO),
    ]
    for i, (name, desc, color) in enumerate(agents):
        left = Inches(0.3) + Inches(i * 3.2)
        _box(slide, f"{name}\n\n{desc}", left, Inches(3.3), Inches(2.9), Inches(1.8),
             fill=color, size=Pt(12), bold=True)
        if i < 3:
            _arrow_r(slide, left + Inches(2.95), Inches(4.0), Inches(0.2), Inches(0.2))

    _box(slide, 'Input: "Mechanisms of antibiotic resistance"  →  Output: Literature report + 3 ranked hypotheses + experiment protocol',
         Inches(0.3), Inches(5.5), Inches(12.6), Inches(0.5), fill=NVIDIA_DARK, color=NVIDIA_WHITE, size=Pt(13))

    _box(slide, "Inspired by Google's AI Co-Scientist  |  Built with Foundation patterns from Labs 0-5  |  Extensible to any domain",
         Inches(0.3), Inches(6.2), Inches(12.6), Inches(0.4), fill=NVIDIA_GRAY, color=NVIDIA_LGRAY, size=Pt(12))

    _notes(slide,
        "This is the flagship lab. Four specialized agents working as a research team, coordinated "
        "by an orchestrator.\n\nLiterature Agent searches papers and finds gaps. Hypothesis Agent "
        "proposes novel ideas. Critic Agent scores each on novelty, feasibility, and impact. "
        "Experiment Agent designs protocols to test the winner.\n\n"
        "One research question goes in, a complete research plan comes out. This is the pattern "
        "behind Google's AI Co-Scientist and NVIDIA's drug discovery multi-agent pipelines.")


def slide_07(prs):
    """Domain highlights — what each track covers."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _common(slide, "DOMAIN TRACKS")
    _title(slide, "What Each Domain Track Teaches")

    _table(slide,
        ["Domain", "Labs", "Key Agent Capabilities", "NVIDIA SDK"],
        [
            ["EOP", "3", "Evidence chain extraction, claim-disclosure mapping, audience advocacy", "NIM"],
            ["Healthcare", "3", "Clinical NLP, PICO literature search, trial eligibility screening", "NIM"],
            ["Bioinformatics", "3", "Sequence analysis, ACMG variant classification, pathway enrichment", "NIM + BioNeMo"],
            ["Finance", "3", "Market analysis + plots, portfolio optimization, algo strategy backtesting", "NIM + cuOpt"],
        ],
        top=Inches(1.7)
    )

    # Key stat boxes
    stats = [
        ("10\ndomain notebooks", NVIDIA_GREEN),
        ("12+ agent tools\nbuilt across tracks", CLR_HEALTH),
        ("15+ visualizations\nwith NVIDIA theme", CLR_BIO),
        ("Every experiment\nhas expected output", CLR_FINANCE),
    ]
    for i, (text, color) in enumerate(stats):
        _box(slide, text, Inches(0.4) + Inches(i * 3.15), Inches(5.0), Inches(2.85), Inches(1.0),
             fill=color, color=NVIDIA_WHITE, size=Pt(14), bold=True)

    _notes(slide,
        "Four domain tracks, each with 3 labs.\n\n"
        "EOP teaches reproducibility — important for any computational research team.\n"
        "Healthcare covers clinical NLP, evidence-based medicine, and trial screening.\n"
        "Bioinformatics goes from DNA sequence to pathway analysis — with a genomics primer for "
        "non-biologists.\n"
        "Finance covers market analysis with 4 matplotlib visualizations, portfolio optimization "
        "with efficient frontier, and algorithmic strategy backtesting.\n\n"
        "Each track builds real tools that scientists can adapt to their own work.")


def slide_08(prs):
    """How to use with customers."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _common(slide, "CUSTOMER PLAYS")
    _title(slide, "How to Use This Playbook with Customers")

    plays = [
        ("Workshop Play", "Half-day guided session through Foundation Labs 0-4.\n"
         "Customer team builds their first agent in 30 minutes.\n"
         "End with Lab 6 AI Co-Scientist demo."),
        ("Domain Play", "Identify customer's field → run the matching domain track.\n"
         "Healthcare customer? HC Labs 1-3 in 3 hours.\n"
         "Finance customer? FIN Labs 1-3 with live portfolio visualizations."),
        ("NIM Conversion Play", "Customer prototypes with OpenAI → show the one-line switch to NIM.\n"
         "Key hook: data privacy + cost at scale + on-prem compliance.\n"
         "Free tier at build.nvidia.com eliminates procurement friction."),
        ("Expansion Play", "Customer completes a domain track → we build a custom domain for them.\n"
         "3-lab template makes this systematic, not ad-hoc.\n"
         "Roadmap domains (BioNeMo, Earth-2, PhysicsNeMo) as upsell."),
    ]
    for i, (title, desc) in enumerate(plays):
        tp = Inches(1.6) + Inches(i * 1.4)
        _box(slide, str(i + 1), Inches(0.4), tp, Inches(0.5), Inches(0.5),
             fill=NVIDIA_GREEN, color=NVIDIA_BLACK, size=Pt(20), bold=True)
        t1 = slide.shapes.add_textbox(Inches(1.2), tp - Inches(0.02), Inches(11.5), Inches(0.3))
        p1 = t1.text_frame.paragraphs[0]
        p1.text = title; p1.font.size = Pt(18); p1.font.color.rgb = NVIDIA_GREEN
        p1.font.bold = True; p1.font.name = FONT
        t2 = slide.shapes.add_textbox(Inches(1.2), tp + Inches(0.35), Inches(11.5), Inches(0.8))
        t2.text_frame.word_wrap = True
        p2 = t2.text_frame.paragraphs[0]
        p2.text = desc; p2.font.size = Pt(12); p2.font.color.rgb = NVIDIA_LGRAY; p2.font.name = FONT

    _notes(slide,
        "Four ways to use this playbook with customers.\n\n"
        "1. Workshop Play: half-day hands-on session. Great for conferences, customer visits, "
        "or onboarding new teams. They leave with a working agent.\n\n"
        "2. Domain Play: match the customer's field to a domain track. Healthcare → HC1-3. "
        "Finance → FIN1-3. They see immediate relevance.\n\n"
        "3. NIM Conversion: they prototype with OpenAI (free, easy), then we show the one-line "
        "switch to NIM. Hook: data privacy, cost, compliance.\n\n"
        "4. Expansion: after they finish a track, we help them build a custom domain. "
        "The 3-lab template makes this repeatable. Roadmap domains are natural upsells to "
        "domain SDK purchases (BioNeMo, Earth-2, etc.).")


def slide_09(prs):
    """What's next — roadmap."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _common(slide, "NEXT STEPS")
    _title(slide, "Roadmap & Action Items")

    items = [
        ("NOW", "Share with your SAs — they can use it in customer engagements this week",
         NVIDIA_GREEN),
        ("THIS MONTH", "Run a pilot workshop with one customer team (healthcare or finance recommended)",
         NVIDIA_GREEN2),
        ("Q2 2026", "Build Drug Discovery domain track with BioNeMo integration",
         RGBColor(0xCC, 0x88, 0x00)),
        ("Q3 2026", "Build Materials + Climate tracks (PhysicsNeMo, Earth-2)",
         RGBColor(0xCC, 0x88, 0x00)),
        ("ONGOING", "Community contributions — customers and partners add their own domain tracks",
         NVIDIA_GRAY),
    ]
    for i, (when, what, color) in enumerate(items):
        tp = Inches(1.6) + Inches(i * 1.1)
        _box(slide, when, Inches(0.4), tp, Inches(2.0), Inches(0.75),
             fill=color, color=NVIDIA_WHITE if color != NVIDIA_GRAY else NVIDIA_LGRAY,
             size=Pt(14), bold=True)
        t = slide.shapes.add_textbox(Inches(2.7), tp + Inches(0.1), Inches(10), Inches(0.5))
        t.text_frame.word_wrap = True
        p = t.text_frame.paragraphs[0]
        p.text = what; p.font.size = Pt(16); p.font.color.rgb = NVIDIA_WHITE; p.font.name = FONT

    _box(slide, "Repo: github.com/NVIDIA/agentic-ai-science-playbook  |  Full customer deck: 64 slides with speaker notes",
         Inches(0.4), Inches(6.4), Inches(12.3), Inches(0.4), fill=NVIDIA_GREEN, color=NVIDIA_BLACK, size=Pt(13), bold=True)

    _notes(slide,
        "Action items for the team:\n\n"
        "Right now: share the playbook repo with your SAs. They can use it in customer meetings "
        "this week — the Scenarios Lab alone is a great conversation starter.\n\n"
        "This month: let's pick one customer and run a pilot workshop. Healthcare and finance "
        "tracks are the most polished — I'd start there.\n\n"
        "Q2-Q3: we build out the NVIDIA SDK domain tracks — drug discovery with BioNeMo, "
        "materials with PhysicsNeMo, climate with Earth-2.\n\n"
        "Long-term: the 3-lab template means customers and partners can contribute their own domains. "
        "This becomes an ecosystem, not just a tutorial.")


def slide_10(prs):
    """Thank you."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(0.15))
    s.fill.solid(); s.fill.fore_color.rgb = NVIDIA_GREEN; s.line.fill.background()

    t = slide.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(3), Inches(0.5))
    p = t.text_frame.paragraphs[0]
    p.text = "NVIDIA"; p.font.size = Pt(18); p.font.color.rgb = NVIDIA_GREEN; p.font.bold = True; p.font.name = FONT

    t2 = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(1))
    p2 = t2.text_frame.paragraphs[0]
    p2.text = "Let's Build This Together"
    p2.font.size = Pt(48); p2.font.color.rgb = NVIDIA_WHITE; p2.font.bold = True; p2.font.name = FONT
    p2.alignment = PP_ALIGN.CENTER

    t3 = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(11), Inches(0.5))
    p3 = t3.text_frame.paragraphs[0]
    p3.text = "Questions?"
    p3.font.size = Pt(28); p3.font.color.rgb = NVIDIA_GREEN; p3.font.name = FONT
    p3.alignment = PP_ALIGN.CENTER

    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(4.6), Inches(4), Inches(0.04))
    ln.fill.solid(); ln.fill.fore_color.rgb = NVIDIA_GREEN; ln.line.fill.background()

    t4 = slide.shapes.add_textbox(Inches(1), Inches(5.0), Inches(11), Inches(1))
    t4.text_frame.word_wrap = True
    for i, txt in enumerate([
        "github.com/NVIDIA/agentic-ai-science-playbook",
        "Full customer deck: 64 slides with speaker notes",
        "build.nvidia.com  |  developer.nvidia.com",
    ]):
        p = t4.text_frame.paragraphs[0] if i == 0 else t4.text_frame.add_paragraph()
        p.text = txt; p.font.size = Pt(14); p.font.color.rgb = NVIDIA_LGRAY; p.font.name = FONT
        p.alignment = PP_ALIGN.CENTER; p.space_after = Pt(4)

    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), SLIDE_H - Inches(0.06), SLIDE_W, Inches(0.06))
    s.fill.solid(); s.fill.fore_color.rgb = NVIDIA_GREEN; s.line.fill.background()

    _notes(slide,
        "Thank you all. I'm excited about this — it's a genuinely useful tool for customer "
        "engagements.\n\nThe full customer-facing deck is 64 slides with speaker notes for every "
        "slide. The 10-pager you just saw is for internal use.\n\n"
        "Who wants to pick the first customer for a pilot workshop?\n\nQuestions?")


# ═══════════════════════════════════════════════════════════════════════════

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_01(prs)  # Title
    slide_02(prs)  # The Problem
    slide_03(prs)  # What's in the playbook
    slide_04(prs)  # Learning roadmap
    slide_05(prs)  # NVIDIA NIM
    slide_06(prs)  # AI Co-Scientist highlight
    slide_07(prs)  # Domain tracks
    slide_08(prs)  # Customer plays
    slide_09(prs)  # Roadmap & action items
    slide_10(prs)  # Thank you

    out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "Playbook_Team_Overview.pptx")
    prs.save(out)
    print(f"Saved {len(prs.slides)} slides -> {out}")


if __name__ == "__main__":
    main()

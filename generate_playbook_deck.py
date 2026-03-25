#!/usr/bin/env python3
"""
Generate an NVIDIA-branded PowerPoint deck introducing the
Agentic AI for Scientific Discovery Playbook.

COMPREHENSIVE EDITION — 50 slides covering:
  - Foundation Labs (7 labs, including AI Co-Scientist)
  - Existing Domain Labs: EOP, Healthcare, Bioinformatics
  - NEW Domain Labs: Drug Discovery, Materials Science, Climate Science,
                     Chemistry, Robotics & Lab Automation
  - NVIDIA SDK Deep Dive (NIM, BioNeMo, Earth-2, PhysicsNeMo, RAPIDS, etc.)
  - Agent Design Principles & Architecture Patterns
  - Production Deployment

Usage:
    pip install python-pptx
    python generate_playbook_deck.py

Output:
    Agentic_AI_Science_Playbook.pptx
"""

from __future__ import annotations

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── NVIDIA Brand Constants ──────────────────────────────────────────────────

NVIDIA_GREEN  = RGBColor(0x76, 0xB9, 0x00)
NVIDIA_BLACK  = RGBColor(0x1A, 0x1A, 0x1A)
NVIDIA_DARK   = RGBColor(0x24, 0x24, 0x24)
NVIDIA_GRAY   = RGBColor(0x2D, 0x2D, 0x2D)
NVIDIA_LGRAY  = RGBColor(0x8C, 0x8C, 0x8C)
NVIDIA_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
NVIDIA_GREEN2 = RGBColor(0x5A, 0x9E, 0x00)
CODE_BG       = RGBColor(0x1E, 0x1E, 0x2E)

# Domain accent colors
CLR_BIO       = RGBColor(0x00, 0x88, 0xCC)  # Blue - bio/genomics
CLR_HEALTH    = RGBColor(0xCC, 0x44, 0x44)  # Red - healthcare
CLR_EOP       = RGBColor(0x88, 0x44, 0xCC)  # Purple - EOP
CLR_DRUG      = RGBColor(0x00, 0xAA, 0x88)  # Teal - drug discovery
CLR_MATERIAL  = RGBColor(0xCC, 0x88, 0x00)  # Amber - materials
CLR_CLIMATE   = RGBColor(0x22, 0x77, 0xBB)  # Sky blue - climate
CLR_CHEM      = RGBColor(0xBB, 0x55, 0x00)  # Orange - chemistry
CLR_ROBOT     = RGBColor(0x66, 0x66, 0xCC)  # Indigo - robotics
CLR_FINANCE   = RGBColor(0x00, 0x88, 0x55)  # Dark green - finance

FONT_TITLE = "Calibri"
FONT_BODY  = "Calibri"
FONT_CODE  = "Consolas"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ── Helpers ─────────────────────────────────────────────────────────────────

def _add_notes(slide, text):
    """Add speaker notes to a slide."""
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = text


def _set_slide_bg(slide, color=NVIDIA_BLACK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_green_bar(slide, top=Inches(0), height=Inches(0.08)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), top, SLIDE_W, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = NVIDIA_GREEN
    shape.line.fill.background()
    return shape


def _add_bottom_bar(slide):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), SLIDE_H - Inches(0.06), SLIDE_W, Inches(0.06))
    shape.fill.solid()
    shape.fill.fore_color.rgb = NVIDIA_GREEN
    shape.line.fill.background()
    return shape


def _add_section_label(slide, text, left=Inches(0.6), top=Inches(0.25)):
    txBox = slide.shapes.add_textbox(left, top, Inches(5), Inches(0.35))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text.upper()
    p.font.size = Pt(11)
    p.font.color.rgb = NVIDIA_GREEN
    p.font.bold = True
    p.font.name = FONT_BODY
    return txBox


def _add_title(slide, text, left=Inches(0.6), top=Inches(0.6), width=Inches(12), size=Pt(36)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = size
    p.font.color.rgb = NVIDIA_WHITE
    p.font.bold = True
    p.font.name = FONT_TITLE
    return txBox


def _add_subtitle(slide, text, left=Inches(0.6), top=Inches(1.3), width=Inches(11), size=Pt(18)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.6))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = size
    p.font.color.rgb = NVIDIA_LGRAY
    p.font.name = FONT_BODY
    return txBox


def _add_bullets(slide, items, left=Inches(0.6), top=Inches(1.9), width=Inches(11.5), size=Pt(17)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (title, desc) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = title
        run.font.size = size
        run.font.color.rgb = NVIDIA_GREEN
        run.font.bold = True
        run.font.name = FONT_BODY
        if desc:
            run2 = p.add_run()
            run2.text = f"  {desc}"
            run2.font.size = Pt(size.pt - 2)
            run2.font.color.rgb = NVIDIA_WHITE
            run2.font.name = FONT_BODY
        p.space_after = Pt(10)
    return txBox


def _add_numbered_list(slide, items, left=Inches(0.6), top=Inches(1.9), width=Inches(11.5), size=Pt(18)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run_num = p.add_run()
        run_num.text = f"{i+1}.  "
        run_num.font.size = size
        run_num.font.color.rgb = NVIDIA_GREEN
        run_num.font.bold = True
        run_num.font.name = FONT_BODY
        run_txt = p.add_run()
        run_txt.text = item
        run_txt.font.size = size
        run_txt.font.color.rgb = NVIDIA_WHITE
        run_txt.font.name = FONT_BODY
        p.space_after = Pt(12)
    return txBox


def _add_code_block(slide, code, left=Inches(0.8), top=Inches(2.0), width=Inches(11.5), height=Inches(4.2)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = CODE_BG
    shape.line.color.rgb = RGBColor(0x3A, 0x3A, 0x4A)
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.25)
    tf.margin_top = Inches(0.2)
    tf.margin_right = Inches(0.25)
    for i, line in enumerate(code.strip().split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(12)
        p.font.color.rgb = NVIDIA_WHITE
        p.font.name = FONT_CODE
        p.space_after = Pt(1)
    return shape


def _add_box(slide, text, left, top, width, height, fill_color=NVIDIA_GRAY, text_color=NVIDIA_WHITE, font_size=Pt(14), bold=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.05)
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        p.text = line
        p.font.size = font_size
        p.font.color.rgb = text_color
        p.font.bold = bold
        p.font.name = FONT_BODY
    return shape


def _add_arrow_right(slide, left, top, width=Inches(0.5), height=Inches(0.3)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = NVIDIA_GREEN
    shape.line.fill.background()
    return shape


def _add_arrow_down(slide, left, top, width=Inches(0.3), height=Inches(0.4)):
    shape = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = NVIDIA_GREEN
    shape.line.fill.background()
    return shape


def _add_table(slide, headers, rows, left=Inches(0.6), top=Inches(2.0), width=Inches(12), row_h=Inches(0.42)):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    col_w = width / n_cols
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, row_h * n_rows)
    table = table_shape.table
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = NVIDIA_BLACK
            p.font.name = FONT_BODY
            p.alignment = PP_ALIGN.LEFT
        cell.fill.solid()
        cell.fill.fore_color.rgb = NVIDIA_GREEN
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12)
                p.font.color.rgb = NVIDIA_WHITE
                p.font.name = FONT_BODY
            cell.fill.solid()
            cell.fill.fore_color.rgb = NVIDIA_DARK if i % 2 == 0 else NVIDIA_GRAY
    for j in range(n_cols):
        table.columns[j].width = int(col_w)
    return table_shape


def _slide_common(slide, section=None):
    _set_slide_bg(slide)
    _add_green_bar(slide)
    _add_bottom_bar(slide)
    if section:
        _add_section_label(slide, section)


def _section_divider(prs, title, subtitle="", notes=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    _add_green_bar(slide)
    _add_bottom_bar(slide)
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11), Inches(1.5))
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44); p.font.color.rgb = NVIDIA_WHITE; p.font.bold = True; p.font.name = FONT_TITLE
    p.alignment = PP_ALIGN.CENTER
    if subtitle:
        txBox2 = slide.shapes.add_textbox(Inches(1), Inches(3.9), Inches(11), Inches(1))
        tf2 = txBox2.text_frame; tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(22); p2.font.color.rgb = NVIDIA_GREEN; p2.font.name = FONT_BODY
        p2.alignment = PP_ALIGN.CENTER
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(3.6), Inches(4), Inches(0.04))
    line.fill.solid(); line.fill.fore_color.rgb = NVIDIA_GREEN; line.line.fill.background()
    if notes:
        _add_notes(slide, notes)
    return slide


def _domain_overview_slide(prs, section_label, title, subtitle, table_rows, pipeline_steps, pipeline_colors, sdk_text, notes=""):
    """Create a standardized domain overview slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_common(slide, section_label)
    _add_title(slide, title)
    _add_subtitle(slide, subtitle)
    _add_table(slide, ["Lab", "Focus", "What You Build"], table_rows, top=Inches(1.8), row_h=Inches(0.4))
    for i, step in enumerate(pipeline_steps):
        left = Inches(0.6) + Inches(i * 2.5)
        color = pipeline_colors[i] if i < len(pipeline_colors) else NVIDIA_GRAY
        _add_box(slide, step, left, Inches(4.6), Inches(2.1), Inches(0.8),
                 fill_color=color, text_color=NVIDIA_WHITE, font_size=Pt(11), bold=True)
        if i < len(pipeline_steps) - 1:
            _add_arrow_right(slide, left + Inches(2.15), Inches(4.85), Inches(0.3), Inches(0.2))
    # SDK badge
    _add_box(slide, sdk_text, Inches(0.6), Inches(5.8), Inches(12), Inches(0.5),
             fill_color=NVIDIA_DARK, text_color=NVIDIA_LGRAY, font_size=Pt(12))
    if notes:
        _add_notes(slide, notes)
    return slide


def _domain_detail_slide(prs, section_label, title, bullets, callout_text, notes=""):
    """Create a standardized domain detail slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_common(slide, section_label)
    _add_title(slide, title)
    _add_bullets(slide, bullets, top=Inches(1.7), size=Pt(16))
    if callout_text:
        _add_box(slide, callout_text, Inches(0.6), Inches(5.6), Inches(12), Inches(0.9),
                 fill_color=NVIDIA_DARK, text_color=NVIDIA_WHITE, font_size=Pt(14))
    if notes:
        _add_notes(slide, notes)
    return slide


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ═══════════════════════════════════════════════════════════════════════════

# ── SECTION 1: TITLE & INTRO (slides 1-3) ──────────────────────────────────

def slide_01_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(0.15))
    bar.fill.solid(); bar.fill.fore_color.rgb = NVIDIA_GREEN; bar.line.fill.background()
    # NVIDIA mark
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(3), Inches(0.5))
    p = txBox.text_frame.paragraphs[0]
    p.text = "NVIDIA"; p.font.size = Pt(18); p.font.color.rgb = NVIDIA_GREEN; p.font.bold = True; p.font.name = FONT_TITLE
    _add_title(slide, "Agentic AI for\nScientific Discovery", top=Inches(1.8), size=Pt(48))
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(3.9), Inches(3), Inches(0.04))
    line.fill.solid(); line.fill.fore_color.rgb = NVIDIA_GREEN; line.line.fill.background()
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(4.2), Inches(11), Inches(0.8))
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "A Comprehensive Hands-On Playbook for Building\nProduction-Quality LLM Agents Across Scientific Domains"
    p.font.size = Pt(22); p.font.color.rgb = NVIDIA_LGRAY; p.font.name = FONT_BODY
    txBox2 = slide.shapes.add_textbox(Inches(0.6), Inches(5.4), Inches(11), Inches(1))
    tf2 = txBox2.text_frame; tf2.word_wrap = True
    for i, d in enumerate([
        "NVIDIA NIM  |  LangGraph  |  Pydantic v2  |  BioNeMo  |  Earth-2  |  PhysicsNeMo",
        "20 Hands-On Labs Today  |  4 Scientific Domains + Finance  |  5 More Domains on Roadmap"
    ]):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.text = d; p.font.size = Pt(14); p.font.color.rgb = NVIDIA_WHITE; p.font.name = FONT_BODY
        p.space_after = Pt(6)
    _add_bottom_bar(slide)
    _add_notes(slide, "Welcome everyone. I'm [name] from NVIDIA, and today I'm going to show you something that's changing how scientists work — AI agents that don't just answer questions, but actually DO research.\n\nThis is a hands-on playbook — 18 labs available today across 4 scientific domains plus finance, with 5 more NVIDIA SDK-powered domains on the roadmap. Your team can start using it today.\n\nThe playbook runs on NVIDIA NIM, but also works with OpenAI — so there's zero barrier to getting started.\n\nLet me show you what's inside.")


def slide_02_agenda(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_common(slide)
    _add_title(slide, "What We'll Cover")
    _add_numbered_list(slide, [
        "Learning Roadmap & 6 Research Scenarios (start here!)",
        "LLM Agent Architecture & Design Principles",
        "Foundation Labs 0-6: Agent Loop to AI Co-Scientist",
        "Domain Labs: EOP, Healthcare, Bioinformatics, Finance",
        "Roadmap: Drug Discovery, Materials, Climate, Chemistry, Robotics (NVIDIA SDK)",
        "NVIDIA NIM Deep Dive: Why NIM for Science?",
        "Agent Patterns, Multi-Agent Orchestration & Production",
        "Getting Started in 5 Minutes",
    ], top=Inches(1.6), size=Pt(17))
    _add_notes(slide, "Here's our roadmap for the next 40 minutes. We'll start with WHY — why should your researchers care about AI agents? Then we'll look at the architecture — what makes an agent different from ChatGPT. Then I'll walk you through the actual labs your team would work through. And we'll end with how to get started — literally today, in 5 minutes.\n\nI want to spend most of our time on the 'so what' — what can these agents actually DO for your team.")


def slide_03_roadmap(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_common(slide, "YOUR LEARNING JOURNEY")
    _add_title(slide, "Learning Roadmap: From Zero to AI Co-Scientist")
    _add_subtitle(slide, "A clear, beginner-friendly path through 18 labs across 5 scientific domains")

    # Roadmap as connected stages
    stages = [
        ("START HERE", "Scenarios Lab\n6 real-world\nuse cases\n(no coding!)", RGBColor(0x44, 0x44, 0xCC), Inches(0.3)),
        ("FOUNDATION", "Labs 0-6\nAgent Loop → Schemas\n→ Memory → Graphs\n→ Eval → Co-Scientist", NVIDIA_GREEN, Inches(2.5)),
        ("CHOOSE YOUR\nDOMAIN", "EOP | Healthcare\nBioinformatics\nFinance\n+ 5 NEW domains", RGBColor(0xCC, 0x88, 0x00), Inches(5.0)),
        ("ADVANCED", "Multi-Agent\nOrchestration\nProduction\nDeployment", NVIDIA_GREEN2, Inches(7.5)),
        ("BUILD YOUR\nOWN", "Apply patterns\nto YOUR research\ndomain", NVIDIA_GREEN, Inches(10.0)),
    ]
    for label, desc, color, left in stages:
        _add_box(slide, f"{label}\n\n{desc}", left, Inches(1.9), Inches(2.2), Inches(3.2),
                 fill_color=color, text_color=NVIDIA_WHITE, font_size=Pt(11), bold=False)
        shape = slide.shapes[-1]
        tf = shape.text_frame
        tf.paragraphs[0].runs[0].font.color.rgb = NVIDIA_WHITE
        tf.paragraphs[0].runs[0].font.bold = True
        tf.paragraphs[0].runs[0].font.size = Pt(14)

    # Arrows between stages
    for left in [Inches(2.55), Inches(5.05), Inches(7.55), Inches(10.05)]:
        _add_arrow_right(slide, left - Inches(0.35), Inches(3.15), Inches(0.3), Inches(0.25))

    # Time estimates
    _add_box(slide, "~45 min  |  ~8 hours  |  ~3 hours per domain  |  ~4 hours  |  Ongoing",
             Inches(0.3), Inches(5.3), Inches(12.3), Inches(0.4),
             fill_color=NVIDIA_DARK, text_color=NVIDIA_LGRAY, font_size=Pt(12))

    # Beginner-friendly features callout
    _add_box(slide,
        "Beginner-Friendly:  Step-by-step code explanations  |  Expected output for every experiment  |  "
        "LangGraph warm-up  |  Genomics primer  |  NIM vs OpenAI decision guide",
        Inches(0.3), Inches(5.9), Inches(12.3), Inches(0.55),
        fill_color=NVIDIA_GRAY, text_color=NVIDIA_WHITE, font_size=Pt(11))

    # Start callout
    _add_box(slide, "New to AI agents? Start with the Scenarios Lab — no coding required, just see what's possible!",
             Inches(0.3), Inches(6.65), Inches(12.3), Inches(0.4),
             fill_color=NVIDIA_GREEN, text_color=NVIDIA_BLACK, font_size=Pt(13), bold=True)
    _add_notes(slide, "This is the learning journey. And I want to highlight something important: you don't have to do all 18 labs. Your team picks a path.\n\nIf they're brand new to agents, they start with the Scenarios Lab — zero coding, just seeing what's possible. That takes 45 minutes.\n\nThen they go through Foundation — that's the core engineering, about 8 hours spread across a few days.\n\nThen they choose THEIR domain — healthcare, genomics, finance, whatever matches their work.\n\nAnd notice the bottom bar: we built this for beginners. Step-by-step explanations, expected output for every experiment, gentle on-ramps for complex topics.\n\nLet me show you what drives all of this.")


def slide_04_scenarios(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_common(slide, "WHY AGENTS?")
    _add_title(slide, "6 Research Scenarios Where AI Agents Help")
    _add_subtitle(slide, "Real problems, real solutions — see what agents can do before learning to build them")

    _add_table(slide,
        ["#", "Scenario", "Pain Point", "What the Agent Does"],
        [
            ["1", "Literature Review", "847 papers to read by Monday", "Search, filter, summarize, find gaps"],
            ["2", "Data Quality Check", "50K rows, something looks wrong", "Scan, flag anomalies, suggest fixes"],
            ["3", "Hypothesis Generation", "200 genes, now what?", "Enrich pathways, propose hypotheses"],
            ["4", "Protocol Writing", "Need a CRISPR protocol ASAP", "Generate steps, safety, materials list"],
            ["5", "Paper Drafting", "6 months of results, blank page", "Outline, draft, cite, format"],
            ["6", "Code Review", "Inherited spaghetti code", "Analyze, document, suggest tests"],
        ],
        top=Inches(1.7)
    )

    _add_box(slide, "Start with the Scenarios Lab — no prior knowledge needed. Then learn to BUILD these agents in the Foundation Labs.",
             Inches(0.6), Inches(5.8), Inches(12), Inches(0.5),
             fill_color=NVIDIA_GREEN, text_color=NVIDIA_BLACK, font_size=Pt(14), bold=True)
    _add_notes(slide, "Before we get technical, let me ask: how many of your researchers have said 'I spend more time wrangling data than doing actual science'?\n\nThese are 6 real scenarios — things that happen every week in research labs. Literature reviews that take weeks. Datasets with hidden quality issues. The blank page when you have results but need to write the paper.\n\nAI agents solve each of these. Not by generating text — by taking ACTION. Searching databases, validating data, designing experiments.\n\nThe Scenarios Lab lets your team experience this in 45 minutes with no coding required. It's the 'movie trailer' before they commit to the full playbook.\n\nNow let me show you how this actually works under the hood.")


def slide_05_landscape(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_common(slide, "THE OPPORTUNITY")
    _add_title(slide, "The Agentic AI Revolution in Science")
    _add_subtitle(slide, "From AI-assisted to AI-driven scientific discovery")
    # Timeline boxes
    eras = [
        ("2020-2022", "AI-Assisted", "Scientists use ML models\nas individual tools", NVIDIA_GRAY),
        ("2023-2024", "AI-Augmented", "LLMs help with literature\nreview and data analysis", NVIDIA_GRAY),
        ("2025-2026", "AI-Agentic", "Autonomous agents plan,\nexecute, and evaluate\nresearch workflows", NVIDIA_GREEN),
        ("2027+", "AI-Collaborative", "Multi-agent teams work\nalongside human scientists\nin closed-loop discovery", NVIDIA_GREEN2),
    ]
    for i, (year, label, desc, color) in enumerate(eras):
        left = Inches(0.4) + Inches(i * 3.2)
        _add_box(slide, f"{year}\n{label}", left, Inches(2.0), Inches(2.8), Inches(0.9),
                 fill_color=color, text_color=NVIDIA_WHITE, font_size=Pt(14), bold=True)
        _add_box(slide, desc, left, Inches(3.1), Inches(2.8), Inches(1.2),
                 fill_color=NVIDIA_DARK, text_color=NVIDIA_WHITE, font_size=Pt(12))
        if i < 3:
            _add_arrow_right(slide, left + Inches(2.85), Inches(2.3), Inches(0.3), Inches(0.25))
    # Stats bar
    _add_box(slide, "80% of scientist time on data wrangling  |  AI agents reduce literature review from weeks to minutes  |  $4.9T healthcare industry adopting 2x faster",
             Inches(0.4), Inches(4.8), Inches(12.4), Inches(0.7),
             fill_color=NVIDIA_DARK, text_color=NVIDIA_LGRAY, font_size=Pt(13))
    _add_notes(slide, "We're at an inflection point. 2020-2022, AI was a tool — you ran a model, got a prediction. 2023-2024, LLMs augmented research — helped with writing and analysis. But NOW, in 2025-2026, we're seeing the shift to AGENTIC AI — systems that plan, execute, and evaluate research workflows autonomously.\n\nThe stat that should get your attention: 80% of a scientist's time goes to data wrangling, not discovery. AI agents flip that ratio.\n\nLet me show you what makes an agent different from a chatbot.")


# ── SECTION 2: AGENT ARCHITECTURE (slides 6-10) ────────────────────────────

def slide_04_what_is_agent(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_common(slide, "AGENT ARCHITECTURE")
    _add_title(slide, "What Is an LLM Agent?")
    _add_subtitle(slide, "Not a chatbot. An autonomous system that reasons, plans, and takes action.")
    boxes = [
        ("PERCEIVE", "Read input, context,\ntool results"),
        ("REASON", "LLM thinks step-by-step\nvia chain-of-thought"),
        ("ACT", "Call tools, execute code,\nquery databases"),
        ("OBSERVE", "Evaluate results,\nupdate memory"),
    ]
    positions = [
        (Inches(1.5), Inches(2.4)), (Inches(7.5), Inches(2.4)),
        (Inches(7.5), Inches(4.6)), (Inches(1.5), Inches(4.6)),
    ]
    bw, bh = Inches(4.0), Inches(1.4)
    for (label, desc), (lft, tp) in zip(boxes, positions):
        box = _add_box(slide, f"{label}\n{desc}", lft, tp, bw, bh,
                       fill_color=NVIDIA_GRAY, text_color=NVIDIA_WHITE, font_size=Pt(13))
        first_p = box.text_frame.paragraphs[0]
        first_p.runs[0].font.color.rgb = NVIDIA_GREEN
        first_p.runs[0].font.bold = True
        first_p.runs[0].font.size = Pt(16)
    _add_arrow_right(slide, Inches(5.6), Inches(2.9), Inches(1.7), Inches(0.3))
    shape = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(9.3), Inches(3.9), Inches(0.3), Inches(0.6))
    shape.fill.solid(); shape.fill.fore_color.rgb = NVIDIA_GREEN; shape.line.fill.background()
    shape = slide.shapes.add_shape(MSO_SHAPE.LEFT_ARROW, Inches(5.6), Inches(5.1), Inches(1.7), Inches(0.3))
    shape.fill.solid(); shape.fill.fore_color.rgb = NVIDIA_GREEN; shape.line.fill.background()
    shape = slide.shapes.add_shape(MSO_SHAPE.UP_ARROW, Inches(3.3), Inches(3.9), Inches(0.3), Inches(0.6))
    shape.fill.solid(); shape.fill.fore_color.rgb = NVIDIA_GREEN; shape.line.fill.background()
    _add_box(slide, "AGENT\nLOOP", Inches(5.6), Inches(3.6), Inches(1.7), Inches(1.0),
             fill_color=NVIDIA_GREEN, text_color=NVIDIA_BLACK, font_size=Pt(16), bold=True)
    _add_notes(slide, "An agent is NOT a chatbot. A chatbot generates text. An agent takes actions.\n\nThis is the core loop: Perceive the problem, Reason about what to do, Act by calling a tool — a database, a simulator, a calculator — then Observe the result and decide what to do next.\n\nThe key insight: the LLM is the brain that decides WHAT to do. Python functions are the hands that DO the work. The LLM never touches your data directly — it orchestrates.\n\nThis separation is what makes agents safe and auditable for science.")


def slide_05_agent_components(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_common(slide, "AGENT ARCHITECTURE")
    _add_title(slide, "Five Pillars of a Scientific Agent")
    components = [
        ("LLM Brain", "Foundation model for\nreasoning & planning\n\nNemotron / GPT-4"),
        ("Tool Library", "Domain-specific functions\nthe agent can invoke\n\nPydantic schemas"),
        ("Memory", "Short-term, long-term,\nepisodic memory\n\nPersists across sessions"),
        ("Orchestration", "State machines for\nmulti-step workflows\n\nLangGraph"),
        ("Evaluation", "Automated quality\nassessment\n\nLLM-as-Judge"),
    ]
    for i, (title, desc) in enumerate(components):
        left = Inches(0.3) + Inches(i * 2.6)
        box = _add_box(slide, f"{title}\n\n{desc}", left, Inches(1.8), Inches(2.3), Inches(4.0),
                       fill_color=NVIDIA_GRAY, text_color=NVIDIA_WHITE, font_size=Pt(12))
        tf = box.text_frame
        tf.paragraphs[0].runs[0].font.color.rgb = NVIDIA_GREEN
        tf.paragraphs[0].runs[0].font.bold = True
        tf.paragraphs[0].runs[0].font.size = Pt(16)
    # Bottom connection bar
    _add_box(slide, "All five pillars are taught progressively across Foundation Labs 0-6",
             Inches(0.3), Inches(6.1), Inches(12.6), Inches(0.5),
             fill_color=NVIDIA_GREEN, text_color=NVIDIA_BLACK, font_size=Pt(14), bold=True)
    _add_notes(slide, "Every scientific agent is built from 5 pillars — and our playbook teaches each one in a dedicated lab.\n\nThe LLM brain for reasoning. A tool library for doing work. Memory so the agent remembers what it learned. Orchestration for multi-step workflows. And evaluation — because in science, if you can't measure it, you can't trust it.\n\nLet me show you why this matters specifically for YOUR researchers.")


def slide_06_why_science(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_common(slide, "AGENT ARCHITECTURE")
    _add_title(slide, "Why LLM Agents for Science?")
    quads = [
        ("Natural Language Interface", "Talk to complex tools in plain English.\nNo API manual needed. Scientists describe\nwhat they want, agents figure out how."),
        ("Structured Reasoning", "Chain-of-thought for hypothesis generation.\nPydantic schemas enforce valid outputs.\nEvery decision is traceable."),
        ("Domain Adaptation", "Same agent patterns work across drug\ndiscovery, climate science, genomics,\nchemistry, materials, and more."),
        ("Composable Workflows", "Chain tools into multi-step pipelines.\nMemory persists. Errors trigger recovery.\nHumans approve critical decisions."),
    ]
    positions = [(Inches(0.5), Inches(1.7)), (Inches(6.7), Inches(1.7)),
                 (Inches(0.5), Inches(4.2)), (Inches(6.7), Inches(4.2))]
    for (title, desc), (lft, tp) in zip(quads, positions):
        box = _add_box(slide, "", lft, tp, Inches(5.9), Inches(2.0), fill_color=NVIDIA_GRAY)
        tf = box.text_frame
        tf.margin_left = Inches(0.2); tf.margin_top = Inches(0.12); tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title; p.font.size = Pt(17); p.font.color.rgb = NVIDIA_GREEN; p.font.bold = True; p.font.name = FONT_BODY; p.alignment = PP_ALIGN.LEFT
        p2 = tf.add_paragraph()
        p2.text = desc; p2.font.size = Pt(13); p2.font.color.rgb = NVIDIA_WHITE; p2.font.name = FONT_BODY; p2.space_before = Pt(6); p2.alignment = PP_ALIGN.LEFT
    _add_notes(slide, "Four reasons agents are a game-changer for science specifically.\n\nFirst: natural language interface. Your biologist doesn't need to learn a PubMed API — they describe what they need in English.\n\nSecond: structured reasoning. Chain-of-thought isn't just a buzzword — it means the agent shows its work, step by step, like a lab notebook.\n\nThird: domain adaptation. The SAME agent pattern works for drug discovery, climate science, genomics. Your team learns once, applies everywhere.\n\nFourth: composable workflows. Chain tools together with memory, error recovery, and human approval gates.\n\nNow let me make the difference concrete.")


def slide_07_chatbot_vs_agent(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_common(slide, "AGENT ARCHITECTURE")
    _add_title(slide, "Chatbot vs. Agent: The Key Differences")
    _add_table(slide,
        ["Capability", "LLM Chatbot", "LLM Agent"],
        [
            ["Interaction", "Single prompt-response", "Multi-step autonomous workflow"],
            ["Tool Use", "None", "Calls functions, APIs, databases, simulators"],
            ["Memory", "Conversation only", "Short-term + long-term + episodic"],
            ["Planning", "None", "Decomposes goals into sub-tasks"],
            ["Error Handling", "Apologizes", "Retries, routes to fallback, escalates"],
            ["Evaluation", "Human reviews output", "LLM-as-Judge with structured rubrics"],
            ["State", "Stateless", "Stateful graph with conditional transitions"],
        ],
        top=Inches(1.7), row_h=Pt(35)
    )
    _add_notes(slide, "This table is the slide I want your team leads to see. Left column: what ChatGPT does. Right column: what an agent does.\n\nLook at error handling: a chatbot apologizes. An agent retries, routes to a fallback, and escalates to a human. That's the difference between a demo and a production system.\n\nAnd evaluation: a chatbot needs a human to review every output. Our agents have built-in LLM-as-Judge scoring — automated quality assessment with structured rubrics.\n\nLet me show you the NVIDIA stack that powers all of this.")


def slide_08_nvidia_stack(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_common(slide, "NVIDIA ECOSYSTEM")
    _add_title(slide, "The NVIDIA AI Stack for Scientific Agents")
    layers = [
        ("NVIDIA NIM  +  Nemotron Models", "GPU-accelerated inference  |  Optimized for reasoning  |  On-prem or cloud", NVIDIA_GREEN),
        ("NeMo Agent Toolkit  +  NeMo Guardrails", "Agent orchestration  |  Safety rails  |  Multi-framework support", RGBColor(0x5A, 0x9E, 0x00)),
        ("Domain SDKs: BioNeMo  |  Earth-2  |  PhysicsNeMo  |  RAPIDS  |  cuOpt", "Pre-trained domain models  |  GPU-accelerated data science  |  Optimization", RGBColor(0x44, 0x77, 0x00)),
        ("LangGraph  +  LangChain  +  Pydantic v2", "State machines  |  Tool chains  |  Type-safe schemas", NVIDIA_GRAY),
        ("CUDA  +  Triton Inference Server  +  Kubernetes", "GPU compute  |  Model serving  |  Production scaling", NVIDIA_DARK),
    ]
    for i, (title, desc, color) in enumerate(layers):
        tp = Inches(1.6) + Inches(i * 1.1)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), tp, Inches(11.5), Inches(0.9))
        box.fill.solid(); box.fill.fore_color.rgb = color; box.line.fill.background()
        tf = box.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.3)
        p = tf.paragraphs[0]
        run = p.add_run(); run.text = title; run.font.size = Pt(16); run.font.bold = True; run.font.color.rgb = NVIDIA_WHITE; run.font.name = FONT_BODY
        p2 = tf.add_paragraph()
        p2.text = desc; p2.font.size = Pt(11); p2.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC); p2.font.name = FONT_BODY
    _add_notes(slide, "This is the full NVIDIA stack for scientific agents, from bottom to top.\n\nAt the foundation: CUDA and Triton for GPU compute. Then LangGraph and Pydantic for orchestration and type safety — open source, framework-agnostic.\n\nIn the middle: BioNeMo for drug discovery, Earth-2 for climate, PhysicsNeMo for materials and chemistry, RAPIDS for data science. These are NVIDIA's domain accelerators.\n\nThen NeMo Agent Toolkit and Guardrails for building safe agents.\n\nAnd at the top: NVIDIA NIM — the inference engine. GPU-optimized, deployable on-prem for data privacy, compatible with the OpenAI API so there's zero migration cost.\n\nLet me show you how the playbook is structured.")


# ── SECTION 3: PLAYBOOK ARCHITECTURE (slides 9-11) ─────────────────────────

def slide_09_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_common(slide, "PLAYBOOK ARCHITECTURE")
    _add_title(slide, "Playbook Architecture: Foundation + Domains")
    _add_subtitle(slide, "18 labs available today + 5 domain expansions on the roadmap")
    # Foundation bar
    _add_box(slide, "FOUNDATION  (7 Labs + Scenarios — Domain-Agnostic Agent Engineering)", Inches(0.5), Inches(1.9), Inches(12.3), Inches(0.55),
             fill_color=NVIDIA_GREEN, text_color=NVIDIA_BLACK, font_size=Pt(15), bold=True)
    # Foundation labs
    labs = ["Lab 0\nPrototype", "Lab 1\nDecisions", "Lab 2\nSchemas", "Lab 3\nMemory", "Lab 4\nGraphs", "Lab 5\nEval", "Lab 6\nCo-Scientist"]
    for i, lab in enumerate(labs):
        _add_box(slide, lab, Inches(0.5) + Inches(i * 1.78), Inches(2.6), Inches(1.58), Inches(0.7),
                 fill_color=NVIDIA_GRAY, text_color=NVIDIA_WHITE, font_size=Pt(10))
    _add_arrow_down(slide, Inches(6.2), Inches(3.4), Inches(0.4), Inches(0.3))
    # AVAILABLE domains - top row
    _add_box(slide, "AVAILABLE TODAY", Inches(0.5), Inches(3.85), Inches(12.3), Inches(0.35),
             fill_color=NVIDIA_GREEN2, text_color=NVIDIA_WHITE, font_size=Pt(12), bold=True)
    avail_domains = [
        ("EOP\n3 Labs", CLR_EOP), ("Healthcare\n3 Labs", CLR_HEALTH),
        ("Bioinformatics\n3 Labs", CLR_BIO), ("Finance\n3 Labs", CLR_FINANCE),
    ]
    for i, (name, color) in enumerate(avail_domains):
        left = Inches(0.5) + Inches(i * 3.15)
        _add_box(slide, name, left, Inches(4.3), Inches(2.85), Inches(0.85),
                 fill_color=color, text_color=NVIDIA_WHITE, font_size=Pt(13), bold=True)
    # ROADMAP domains - bottom row
    _add_box(slide, "ON ROADMAP — NVIDIA SDK Integration", Inches(0.5), Inches(5.35), Inches(12.3), Inches(0.35),
             fill_color=NVIDIA_DARK, text_color=NVIDIA_LGRAY, font_size=Pt(12), bold=True)
    road_domains = [
        ("Drug Discovery\nBioNeMo", CLR_DRUG), ("Materials\nPhysicsNeMo", CLR_MATERIAL),
        ("Climate\nEarth-2", CLR_CLIMATE), ("Chemistry\ncuOpt + Modulus", CLR_CHEM),
        ("Robotics\nIsaac Sim", CLR_ROBOT),
    ]
    for i, (name, color) in enumerate(road_domains):
        left = Inches(0.5) + Inches(i * 2.5)
        _add_box(slide, name, left, Inches(5.8), Inches(2.25), Inches(0.85),
                 fill_color=NVIDIA_GRAY, text_color=NVIDIA_LGRAY, font_size=Pt(11), bold=True)
    _add_notes(slide, "Two layers. The foundation — 7 labs plus a Scenarios intro, all domain-agnostic.\n\nThen domains. TODAY we have 4 fully built domain tracks: EOP with 3 labs, Healthcare with 3 labs, Bioinformatics with 3 labs, and Finance with 1 lab. That's 18 hands-on notebooks your team can use right now.\n\nOn the roadmap: Drug Discovery with BioNeMo, Materials with PhysicsNeMo, Climate with Earth-2, Chemistry with cuOpt, and Robotics with Isaac Sim. Each will follow the same 3-lab template.\n\nThe key insight: the foundation patterns TRANSFER. Once your team masters the 7 foundation labs, they can apply those skills to ANY domain.")


def slide_10_learning_paths(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_common(slide, "PLAYBOOK ARCHITECTURE")
    _add_title(slide, "Learning Paths — Pick Your Journey")
    # AVAILABLE paths
    paths = [
        ("Path A", "Agent Fundamentals", "Foundation Labs 0-6", NVIDIA_GREEN, "Any domain scientist"),
        ("Path B", "Computational Science", "Labs 0-4 + EOP 1-3", CLR_EOP, "PhD students, research engineers"),
        ("Path C", "Healthcare AI", "Labs 0-4 + HC 1-3", CLR_HEALTH, "Clinical informatics, health tech"),
        ("Path D", "Bioinformatics", "Labs 0-4 + BIO 1-3", CLR_BIO, "Genomics, proteomics researchers"),
        ("Path E", "Quantitative Finance", "Labs 0-4 + FIN 1-3", CLR_FINANCE, "Quants, risk analysts, fintech"),
        ("Path F", "Full Mastery", "All 20 Labs", RGBColor(0xCC, 0x88, 0x00), "AI platform architects"),
    ]
    _add_box(slide, "AVAILABLE TODAY", Inches(0.4), Inches(1.5), Inches(12.3), Inches(0.3),
             fill_color=NVIDIA_GREEN2, text_color=NVIDIA_WHITE, font_size=Pt(11), bold=True)
    for i, (label, name, labs, color, audience) in enumerate(paths):
        tp = Inches(1.95) + Inches(i * 0.62)
        _add_box(slide, label, Inches(0.4), tp, Inches(1.2), Inches(0.5),
                 fill_color=color, text_color=NVIDIA_WHITE, font_size=Pt(13), bold=True)
        txBox = slide.shapes.add_textbox(Inches(1.8), tp + Inches(0.02), Inches(5.5), Inches(0.3))
        p = txBox.text_frame.paragraphs[0]
        p.text = f"{name}  ({labs})"; p.font.size = Pt(14); p.font.color.rgb = NVIDIA_WHITE; p.font.bold = True; p.font.name = FONT_BODY
        txBox2 = slide.shapes.add_textbox(Inches(8.5), tp + Inches(0.05), Inches(4.5), Inches(0.3))
        p2 = txBox2.text_frame.paragraphs[0]
        p2.text = audience; p2.font.size = Pt(12); p2.font.color.rgb = NVIDIA_LGRAY; p2.font.name = FONT_BODY
    # ROADMAP paths
    _add_box(slide, "ON ROADMAP", Inches(0.4), Inches(5.7), Inches(12.3), Inches(0.3),
             fill_color=NVIDIA_DARK, text_color=NVIDIA_LGRAY, font_size=Pt(11), bold=True)
    road_paths = [
        ("Drug Discovery (BioNeMo)", CLR_DRUG), ("Materials (PhysicsNeMo)", CLR_MATERIAL),
        ("Climate (Earth-2)", CLR_CLIMATE), ("Chemistry (cuOpt)", CLR_CHEM), ("Robotics (Isaac)", CLR_ROBOT),
    ]
    for i, (name, color) in enumerate(road_paths):
        left = Inches(0.4) + Inches(i * 2.55)
        _add_box(slide, name, left, Inches(6.15), Inches(2.35), Inches(0.55),
                 fill_color=NVIDIA_GRAY, text_color=NVIDIA_LGRAY, font_size=Pt(10), bold=True)
    _add_notes(slide, "Six paths available today. Path A is pure agent engineering — any domain scientist. Paths B-E are domain tracks with 3 labs each. Path F is the full mastery track for platform architects.\n\nOn the roadmap: 5 more domain tracks integrating NVIDIA's domain SDKs — BioNeMo for drug discovery, PhysicsNeMo for materials, Earth-2 for climate, cuOpt for chemistry, and Isaac Sim for robotics.\n\nWhich paths match your team? [PAUSE for customer input]")


def slide_11_tech_stack(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_common(slide, "PLAYBOOK ARCHITECTURE")
    _add_title(slide, "Tech Stack & Configuration")
    _add_table(slide,
        ["Component", "Choice", "Why"],
        [
            ["LLM Provider", "OpenAI / NVIDIA NIM", "One env var to switch — same code"],
            ["Agent Framework", "Raw API (0-3), LangGraph (4+)", "Gradual complexity increase"],
            ["Schemas", "Pydantic v2", "Industry-standard type-safe output"],
            ["Evaluation", "LLM-as-Judge", "Scalable automated quality signal"],
            ["Domain SDKs", "BioNeMo, Earth-2, PhysicsNeMo", "GPU-accelerated domain models"],
            ["Notebooks", "Jupyter / Google Colab", "Zero-friction reproducibility"],
        ],
        top=Inches(1.6), row_h=Pt(32)
    )
    _add_code_block(slide,
        '# Switch to NVIDIA NIM — one line\n'
        'export USE_NIM=true\n'
        'export NIM_API_KEY="nvapi-..."           # from build.nvidia.com\n'
        '\n'
        'client, model = make_client()             # Auto-configures\n'
        '# Provider: NVIDIA NIM\n'
        '# Model:    nvidia/llama-3.3-nemotron-super-49b-v1.5',
        top=Inches(5.0), height=Inches(1.9)
    )
    _add_notes(slide, "The tech stack is deliberately simple. Python, OpenAI or NIM, Pydantic for schemas, LangGraph for orchestration.\n\nAnd here's the key: switching from OpenAI to NVIDIA NIM is one environment variable. Same code, same tools, same schemas. Your team prototypes with OpenAI, deploys with NIM for data privacy and cost control.\n\nLet me walk you through the actual labs.")


# ── SECTION 4: FOUNDATION LABS (slides 12-18) ──────────────────────────────

def slide_12_foundation_overview(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_common(slide, "FOUNDATION LABS")
    _add_title(slide, "Foundation Labs: 7 Labs, 7 Production Patterns")
    _add_subtitle(slide, "Domain-agnostic skills every scientific agent needs")
    _add_table(slide,
        ["Lab", "Core Skill", "Key Concept", "Pattern"],
        [
            ["Lab 0", "Agent Prototype", "Prompt -> LLM -> parse -> execute", "Minimal agent loop"],
            ["Lab 1", "Decision Anatomy", "Prompt structure -> tool selection", "Prompt engineering"],
            ["Lab 2", "Tool Contracts", "Pydantic schemas -> function calls", "Type-safe interfaces"],
            ["Lab 3", "Persistent Agent", "3-layer memory architecture", "State management"],
            ["Lab 4", "Graphs & Recovery", "LangGraph state machines", "Orchestration"],
            ["Lab 5", "LLM-as-Judge", "Multi-dimensional rubrics", "Automated evaluation"],
            ["Lab 6", "AI Co-Scientist", "Multi-agent research team", "Multi-agent collaboration"],
        ],
        top=Inches(1.7)
    )
    _add_notes(slide, "7 labs, 7 production patterns. Each one takes about an hour.\n\nNotice the progression: Lab 0 builds a working agent in 50 lines of Python. Lab 1 teaches you why prompts matter. Lab 2 adds type safety. Lab 3 adds memory. Lab 4 adds orchestration. Lab 5 adds evaluation. And Lab 6 — our newest — builds a multi-agent AI Co-Scientist.\n\nLet me show you the code for Lab 0 — it's simpler than you'd expect.")


def slide_13_lab0(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_common(slide, "FOUNDATION — LAB 0")
    _add_title(slide, "Lab 0: The Minimal Agent Loop")
    _add_subtitle(slide, "Build an agent from scratch in ~50 lines of Python")
    _add_code_block(slide,
        'def agent_loop(user_input: str, tools: dict, max_turns=5):\n'
        '    messages = [{"role": "system", "content": SYSTEM_PROMPT}]\n'
        '    messages.append({"role": "user", "content": user_input})\n'
        '\n'
        '    for turn in range(max_turns):\n'
        '        response = client.chat.completions.create(\n'
        '            model=model, messages=messages\n'
        '        )\n'
        '        reply = response.choices[0].message.content\n'
        '\n'
        '        tool_name = parse_tool_choice(reply)    # TOOL: <name>\n'
        '        if tool_name is None:\n'
        '            return reply                         # Final answer\n'
        '\n'
        '        result = tools[tool_name](reply)         # Execute tool\n'
        '        messages.append({"role": "assistant", "content": reply})\n'
        '        messages.append({"role": "user",\n'
        '                         "content": f"Tool result: {result}"})\n'
        '\n'
        '    return "Max turns reached"',
        top=Inches(1.6), height=Inches(5.0)
    )
    _add_notes(slide, "This is the entire agent. 20 lines of Python. It takes a user message, builds a prompt, calls the LLM, parses which tool to use, executes it, and returns the result.\n\nEvery agent — from a simple research assistant to a multi-agent drug discovery pipeline — is built on this loop.\n\nYour team can run this in a Jupyter notebook in under 10 minutes. First lab, first working agent.\n\nNow, this is simple by design. The next labs add the sophistication.")


def slide_14_lab1_lab2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_common(slide, "FOUNDATION — LABS 1 & 2")
    _add_title(slide, "Labs 1-2: Decisions & Tool Contracts")
    # Left: Lab 1
    _add_box(slide, "Lab 1: Anatomy of a Decision", Inches(0.4), Inches(1.6), Inches(6.1), Inches(0.5),
             fill_color=NVIDIA_GREEN2, text_color=NVIDIA_WHITE, font_size=Pt(15), bold=True)
    _add_bullets(slide, [
        ("Prompt structure", "controls which tool the agent selects"),
        ("Temperature", "low = deterministic, high = creative but risky"),
        ("Experiment design", "run same query across prompt variants"),
    ], left=Inches(0.5), top=Inches(2.3), width=Inches(5.8), size=Pt(14))
    # Right: Lab 2
    _add_box(slide, "Lab 2: Contract of a Tool", Inches(6.8), Inches(1.6), Inches(6.1), Inches(0.5),
             fill_color=NVIDIA_GREEN2, text_color=NVIDIA_WHITE, font_size=Pt(15), bold=True)
    _add_code_block(slide,
        'class SearchPubMed(BaseModel):\n'
        '    """Search PubMed."""\n'
        '    query: str = Field(\n'
        '        description="Search terms"\n'
        '    )\n'
        '    max_results: int = Field(\n'
        '        default=5, ge=1, le=20\n'
        '    )\n'
        '\n'
        '# Auto-generates JSON schema\n'
        '# -> OpenAI tools format',
        left=Inches(6.9), top=Inches(2.3), width=Inches(5.8), height=Inches(3.8)
    )
    _add_box(slide, '"The prompt IS the architecture"  |  "If your tool doesn\'t have a schema, your agent is guessing"',
             Inches(0.4), Inches(6.3), Inches(12.4), Inches(0.5),
             fill_color=NVIDIA_DARK, text_color=NVIDIA_GREEN, font_size=Pt(13), bold=True)
    _add_notes(slide, "Two crucial upgrades. Lab 1 shows that prompt structure DRIVES agent behavior — small changes in wording cause large changes in which tool gets selected. We run controlled experiments, like a scientist would.\n\nLab 2 replaces our fragile regex parser with Pydantic schemas and OpenAI function calling. Now every tool argument is typed, validated, and documented. This is the difference between a demo and production.\n\nThe key quote: 'If your tool doesn't have a schema, your agent is guessing.'")


def slide_15_lab3(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_common(slide, "FOUNDATION — LAB 3")
    _add_title(slide, "Lab 3: Three-Layer Memory Architecture")
    _add_subtitle(slide, "Memory makes agents learn and adapt across sessions")
    layers = [
        ("SHORT-TERM", "Conversation History", "Recent messages — forgotten when session ends", NVIDIA_GREEN),
        ("LONG-TERM", "Key-Value Knowledge Store", "Facts, preferences, domain knowledge — persists across sessions", RGBColor(0x4A, 0x7A, 0x00)),
        ("EPISODIC", "Action Log", "What the agent did, when, why — enables reflection and learning", NVIDIA_GRAY),
    ]
    for i, (label, title, desc, color) in enumerate(layers):
        tp = Inches(2.0) + Inches(i * 1.5)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), tp, Inches(11), Inches(1.2))
        box.fill.solid(); box.fill.fore_color.rgb = color; box.line.fill.background()
        tf = box.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.3); tf.margin_top = Inches(0.08)
        p = tf.paragraphs[0]
        run = p.add_run(); run.text = f"{label}  —  {title}"; run.font.size = Pt(17); run.font.bold = True; run.font.color.rgb = NVIDIA_WHITE; run.font.name = FONT_BODY
        p2 = tf.add_paragraph()
        p2.text = desc; p2.font.size = Pt(13); p2.font.color.rgb = RGBColor(0xDD, 0xDD, 0xDD); p2.font.name = FONT_BODY; p2.space_before = Pt(2)
    for i in range(2):
        _add_arrow_down(slide, Inches(6.3), Inches(3.25) + Inches(i * 1.5), Inches(0.35), Inches(0.25))
    _add_notes(slide, "Memory is what separates a chatbot from a research assistant. Three layers:\n\nShort-term: the current conversation. Long-term: facts that persist — 'this user works on CRISPR, prefers bullet points.' Episodic: what the agent did and when — a digital lab notebook.\n\nFor your researchers, this means the agent remembers last week's findings and builds on them. No more re-explaining context every session.")


def slide_16_lab4(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_common(slide, "FOUNDATION — LAB 4")
    _add_title(slide, "Lab 4: LangGraph State Machines")
    _add_subtitle(slide, '"Real agents fail gracefully" — conditional edges for error recovery')
    nodes = [
        ("START", Inches(0.6), Inches(2.5), NVIDIA_GREEN),
        ("CLASSIFY", Inches(2.8), Inches(2.5), NVIDIA_GRAY),
        ("EXECUTE\nTOOL", Inches(5.5), Inches(2.5), NVIDIA_GRAY),
        ("VALIDATE", Inches(8.2), Inches(2.5), NVIDIA_GRAY),
        ("DONE", Inches(11.0), Inches(2.5), NVIDIA_GREEN),
        ("RETRY", Inches(5.5), Inches(4.6), RGBColor(0x8B, 0x00, 0x00)),
        ("HUMAN\nREVIEW", Inches(8.2), Inches(4.6), RGBColor(0xCC, 0x88, 0x00)),
    ]
    for name, lft, tp, color in nodes:
        _add_box(slide, name, lft, tp, Inches(1.6), Inches(0.75),
                 fill_color=color, text_color=NVIDIA_WHITE, font_size=Pt(12), bold=True)
    for lft in [Inches(2.2), Inches(4.45), Inches(7.15), Inches(9.85)]:
        _add_arrow_right(slide, lft, Inches(2.7), Inches(0.5), Inches(0.25))
    _add_arrow_down(slide, Inches(6.1), Inches(3.3), Inches(0.25), Inches(1.2))
    shape = slide.shapes.add_shape(MSO_SHAPE.UP_ARROW, Inches(5.65), Inches(3.35), Inches(0.25), Inches(1.15))
    shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(0x8B, 0x00, 0x00); shape.line.fill.background()
    _add_arrow_down(slide, Inches(8.8), Inches(3.3), Inches(0.25), Inches(1.2))
    _add_subtitle(slide, "Conditional edges let agents retry failed steps, escalate to human review, or route to alternative tools",
                  top=Inches(5.8), size=Pt(13))
    _add_notes(slide, "Real research workflows aren't linear. Searches fail. Simulations time out. Results need human review.\n\nLangGraph lets us model this as a state machine. Look at the retry cycle: if search returns nothing, the agent automatically refines the query and tries again. If it fails 3 times, it escalates to human review.\n\nWe added a warm-up section — your team starts with a 2-node graph before the full pipeline. No steep learning curve.")


def slide_17_lab5(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_common(slide, "FOUNDATION — LAB 5")
    _add_title(slide, "Lab 5: LLM-as-Judge Evaluation")
    _add_subtitle(slide, '"If you can\'t measure it, you can\'t improve it"')
    _add_table(slide,
        ["Dimension", "Scale", "What It Measures"],
        [
            ["Correctness", "1-5", "Are the facts and tool calls accurate?"],
            ["Completeness", "1-5", "Did the agent address all aspects?"],
            ["Groundedness", "1-5", "Are claims backed by evidence?"],
            ["Clarity", "1-5", "Is the output clear and well-structured?"],
        ],
        top=Inches(1.7)
    )
    _add_code_block(slide,
        'class EvalRubric(BaseModel):\n'
        '    correctness:   int = Field(ge=1, le=5)\n'
        '    completeness:  int = Field(ge=1, le=5)\n'
        '    groundedness:  int = Field(ge=1, le=5)\n'
        '    clarity:       int = Field(ge=1, le=5)\n'
        '    justification: str\n'
        '\n'
        'score = llm_judge(agent_output, rubric=EvalRubric)',
        top=Inches(4.2), height=Inches(2.4)
    )
    _add_notes(slide, "In science, you don't trust results without metrics. Same for agents.\n\nLLM-as-Judge scores every agent output on 4 dimensions: correctness, completeness, groundedness, and clarity. 1-5 scale, Pydantic-validated rubric.\n\nThis enables A/B testing: which prompt is better? Which model produces more grounded responses? Your team can answer these objectively, not anecdotally.")


def slide_18_foundation_summary(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_common(slide, "FOUNDATION LABS")
    _add_title(slide, "7 Patterns You'll Master")
    patterns = [
        ("1. Minimal Agent Loop", "Prompt -> LLM -> parse\n-> execute -> repeat"),
        ("2. Prompt-Driven Decisions", "Prompt structure controls\ntool selection probability"),
        ("3. Type-Safe Contracts", "Pydantic schemas enforce\nvalidated arguments"),
        ("4. Three-Layer Memory", "Short + long + episodic\nmemory architecture"),
        ("5. Graph Orchestration", "LangGraph state machines\nwith error recovery"),
        ("6. Automated Evaluation", "LLM-as-Judge with\nmulti-dimensional rubrics"),
        ("7. Multi-Agent Research", "AI Co-Scientist with\nspecialized agent teams"),
    ]
    for i, (title, desc) in enumerate(patterns):
        col, row = i % 4, i // 4
        left = Inches(0.3) + Inches(col * 3.2)
        top = Inches(1.7) + Inches(row * 2.7)
        box = _add_box(slide, f"{title}\n\n{desc}", left, top, Inches(3.0), Inches(2.3),
                       fill_color=NVIDIA_GRAY, text_color=NVIDIA_WHITE, font_size=Pt(12))
        tf = box.text_frame
        tf.paragraphs[0].runs[0].font.color.rgb = NVIDIA_GREEN
        tf.paragraphs[0].runs[0].font.bold = True
        tf.paragraphs[0].runs[0].font.size = Pt(14)
    _add_notes(slide, "7 patterns, each one production-grade. Your team comes out of the foundation with a toolkit that works in any scientific domain.\n\nBut the newest addition is my favorite...")


def slide_19_coscientist(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_common(slide, "FOUNDATION — LAB 6")
    _add_title(slide, "Lab 6: AI Co-Scientist — Multi-Agent Research")
    _add_subtitle(slide, "Four specialized agents collaborate on scientific discovery")

    # Agent boxes in pipeline
    agents = [
        ("LITERATURE\nAGENT", "Searches papers\nFinds gaps", CLR_HEALTH),
        ("HYPOTHESIS\nAGENT", "Generates novel\nideas from gaps", CLR_DRUG),
        ("CRITIC\nAGENT", "Scores novelty\nfeasibility, impact", RGBColor(0xCC, 0x88, 0x00)),
        ("EXPERIMENT\nAGENT", "Designs protocols\nto test hypotheses", CLR_BIO),
    ]

    # Orchestrator at top
    _add_box(slide, "ORCHESTRATOR\nCoordinates the research team", Inches(4.2), Inches(1.7), Inches(4.8), Inches(0.8),
             fill_color=NVIDIA_GREEN, text_color=NVIDIA_BLACK, font_size=Pt(13), bold=True)

    # Arrow down from orchestrator
    _add_arrow_down(slide, Inches(6.4), Inches(2.55), Inches(0.3), Inches(0.4))

    # 4 agents in a row
    for i, (name, desc, color) in enumerate(agents):
        left = Inches(0.3) + Inches(i * 3.2)
        _add_box(slide, f"{name}\n\n{desc}", left, Inches(3.2), Inches(2.9), Inches(1.8),
                 fill_color=color, text_color=NVIDIA_WHITE, font_size=Pt(12), bold=True)
        if i < 3:
            _add_arrow_right(slide, left + Inches(2.95), Inches(3.9), Inches(0.2), Inches(0.2))

    _add_box(slide, "Inspired by Google's AI Co-Scientist  |  Uses Pydantic structured output  |  Each agent is a focused LLM call  |  Orchestrator aggregates results",
             Inches(0.3), Inches(5.4), Inches(12.6), Inches(0.5),
             fill_color=NVIDIA_DARK, text_color=NVIDIA_LGRAY, font_size=Pt(12))

    # Example output
    _add_box(slide, 'Example: "Mechanisms of antibiotic resistance"\n→ Literature: 5 papers, 3 gaps  →  Hypotheses: 3 ranked ideas  →  Critique: top pick scored 8.5/10  →  Experiment: RCT protocol designed',
             Inches(0.3), Inches(6.1), Inches(12.6), Inches(0.5),
             fill_color=NVIDIA_GRAY, text_color=NVIDIA_WHITE, font_size=Pt(11))
    _add_notes(slide, "Lab 6: the AI Co-Scientist. Four specialized agents working as a research team.\n\nThe Literature Agent searches papers and finds gaps. The Hypothesis Agent proposes novel ideas from those gaps. The Critic scores each hypothesis on novelty, feasibility, and impact. The Experiment Agent designs protocols to test the winner.\n\nThis is inspired by Google's AI Co-Scientist — and your team builds their own version in this lab.\n\nThe orchestrator at the top coordinates everything. One research question goes in, a full research plan comes out.\n\nNow let's see the domain tracks.")


# ── SECTION 5: EXISTING DOMAIN LABS ────────────────────────────────────────

def slide_20_eop(prs):
    _domain_overview_slide(prs,
        "DOMAIN — EVIDENCE-ORIENTED PROGRAMMING", "EOP: Making Research Reproducible",
        "Traceable research software for computational science",
        [
            ["EOP 1", "Evidence Chain Extraction", "Identify ECF's 7 artifact types in research repos"],
            ["EOP 2", "Claim-Contingent Disclosure", "Map claim strength to required disclosure scope"],
            ["EOP 3", "EOP Spokesperson", "Advocate EOP to different audiences with LLM-as-Judge"],
        ],
        ["Input\nData", "Experimental\nProcess", "Output", "Visual\nClaim", "Documentation"],
        [CLR_EOP, NVIDIA_GRAY, CLR_EOP, NVIDIA_GRAY, CLR_EOP],
        "NVIDIA SDKs: NIM (Nemotron)  |  Pattern: Literature Review + Hypothesis Generation  |  Framework: Zhang et al. ECF",
        notes="Evidence-Oriented Programming ensures that computational research is reproducible. The agent analyzes research repos, classifies evidence artifacts, and validates that every claim has complete support.\n\nThis matters because 70% of computational research can't be reproduced. EOP agents can audit a repo in seconds."
    )


def slide_21_eop_detail(prs):
    _domain_detail_slide(prs,
        "DOMAIN — EOP", "EOP Deep Dive: Claim Strength & Disclosure",
        [
            ("Existential Claims", '"Feature X exists" -> Level 1: Code availability'),
            ("Comparative Claims", '"Method A outperforms B" -> Level 2: Code + data + parameters'),
            ("Distributional Claims", '"Effect holds across conditions" -> Level 3: Full pipeline + configs'),
            ("Novel Method Claims", '"We introduce a new algorithm" -> Level 4: Complete reproduction package'),
        ],
        "Agent reads a paper, classifies each claim by strength, determines minimum disclosure for reproducibility.",
        notes="The key insight: stronger claims need stronger evidence. A simple 'feature X exists' needs basic code availability. But 'our method outperforms all baselines' needs the full pipeline, data, parameters, and configs.\n\nThe agent classifies each claim and tells you exactly what you need to disclose."
    )


def slide_22_healthcare(prs):
    _domain_overview_slide(prs,
        "DOMAIN — HEALTHCARE", "Healthcare: From Clinical Notes to Evidence-Based Decisions",
        "Clinical NLP, medical literature synthesis, and trial assistance",
        [
            ["HC 1", "Clinical NLP Agent", "Extract medications, diagnoses, vitals; detect allergy conflicts"],
            ["HC 2", "Medical Literature Agent", "PICO search + GRADE evidence quality + multi-study synthesis"],
            ["HC 3", "Clinical Trial Assistant", "Patient eligibility screening against inclusion/exclusion criteria"],
        ],
        ["Clinical\nNotes", "NLP\nExtraction", "Literature\nSearch", "Evidence\nGrading", "Decision\nSupport"],
        [NVIDIA_GRAY, CLR_HEALTH, NVIDIA_GRAY, CLR_HEALTH, NVIDIA_GRAY],
        "NVIDIA SDKs: NIM (Nemotron)  |  Patterns: Data Analysis + Literature Review  |  Standards: PICO, GRADE, HL7",
        notes="Healthcare generates 30% of the world's data, mostly unstructured clinical notes. These three labs build agents that extract structured data, search medical literature using PICO, and screen patients for clinical trials.\n\nIQVIA's partnership with NVIDIA already cut clinical data review from 7 weeks to 2 weeks using exactly these patterns."
    )


def slide_23_healthcare_detail(prs):
    _domain_detail_slide(prs,
        "DOMAIN — HEALTHCARE", "Healthcare Deep Dive: Medical Literature Agent",
        [
            ("PICO Framework", "Population, Intervention, Comparison, Outcome — structured search queries"),
            ("GRADE Evidence Pyramid", "Agent assesses quality: High / Moderate / Low / Very Low"),
            ("Multi-Study Synthesis", "Synthesizes findings across papers into actionable clinical guidance"),
            ("Allergy Conflict Detection", "Flags drug interactions and contraindications automatically"),
            ("Protocol Explanation", "Translates complex trial protocols into patient-friendly language"),
        ],
        "HC2 demonstrates a complete agentic RAG pipeline: search -> retrieve -> evaluate -> synthesize -> present.",
        notes="The Medical Literature Agent demonstrates the full agentic RAG pipeline. PICO-structured search, GRADE evidence assessment, multi-study synthesis.\n\nThis is a systematic review in seconds, not months. And every source is cited and graded."
    )


def slide_24_bioinformatics(prs):
    _domain_overview_slide(prs,
        "DOMAIN — BIOINFORMATICS", "Bioinformatics: From DNA Sequence to Biological Insight",
        "Sequence analysis, variant interpretation, and pathway reasoning",
        [
            ["BIO 1", "Sequence Analysis Agent", "GC content, ORF detection, motif search, DNA-to-protein translation"],
            ["BIO 2", "Variant Interpretation", "ACMG/AMP classification with pathogenicity assessment"],
            ["BIO 3", "Pathway Analysis Agent", "GO/KEGG enrichment, drug targets, interaction networks"],
        ],
        ["DNA/RNA\nSequence", "Sequence\nAnalysis", "Variant\nCalling", "Pathogenicity", "Pathway\nEnrichment"],
        [NVIDIA_GRAY, CLR_BIO, NVIDIA_GRAY, CLR_BIO, NVIDIA_GRAY],
        "NVIDIA SDKs: NIM, BioNeMo (ESMFold, MolMIM)  |  Patterns: Data Analysis + Hypothesis Generation  |  Standards: ACMG/AMP",
        notes="From DNA sequence to biological insight. Sequence analysis, variant interpretation, pathway enrichment. Each lab connects to NVIDIA's genomics accelerators — BioNeMo for protein structure, Parabricks for variant calling.\n\nAnd we added a genomics primer for non-biologists — your AI engineers can follow this without a biology degree."
    )


def slide_25_bioinformatics_detail(prs):
    _domain_detail_slide(prs,
        "DOMAIN — BIOINFORMATICS", "Bioinformatics Deep Dive: Variant Interpretation",
        [
            ("ACMG/AMP Classification", "5-tier system: Pathogenic / Likely Pathogenic / VUS / Likely Benign / Benign"),
            ("Evidence Criteria", "Agent evaluates population frequency, computational predictions, functional data"),
            ("Structured Report", "Pydantic schema enforces standardized pathogenicity output"),
            ("Gene-Disease Associations", "Cross-references ClinVar, gnomAD, OMIM databases"),
            ("Actionable Insights", "Clinical recommendations based on variant classification"),
        ],
        "Agent receives a variant (e.g., BRCA1 c.5266dupC), queries databases, applies ACMG criteria, produces structured report.",
        notes="ACMG variant classification is the gold standard for clinical genomics. The agent queries databases, evaluates evidence, and produces structured pathogenicity reports.\n\nWe walk through the full notation system — even if you've never seen 'BRCA1:c.5266dupC' before, the primer explains it in 2 minutes."
    )


# ── SECTION 6: NEW DOMAIN LABS (slides 25-34) ──────────────────────────────

def slide_25_drug_discovery(prs):
    _domain_overview_slide(prs,
        "ROADMAP — DRUG DISCOVERY", "Drug Discovery: AI-Driven Molecular Design",
        "Accelerate the drug discovery pipeline with LLM agents and BioNeMo",
        [
            ["DRUG 1", "Molecular Generation Agent", "Generate novel candidate molecules with desired properties"],
            ["DRUG 2", "Docking & Scoring Agent", "Predict binding affinity and rank compounds"],
            ["DRUG 3", "Multi-Agent Drug Pipeline", "End-to-end: target ID -> generation -> scoring -> ADMET"],
        ],
        ["Target\nID", "Molecule\nGeneration", "Docking &\nScoring", "ADMET\nPrediction", "Lead\nOptimization"],
        [NVIDIA_GRAY, CLR_DRUG, CLR_DRUG, NVIDIA_GRAY, CLR_DRUG],
        "NVIDIA SDKs: BioNeMo (MolMIM, DiffDock, ESM-2), NIM, RAPIDS  |  Patterns: Simulation + Multi-Agent Orchestration",
        notes="Drug discovery is the highest-impact use case. BioNeMo provides MolMIM for molecular generation, DiffDock for docking, ESM-2 for protein analysis.\n\nThe multi-agent pipeline goes from target identification through lead optimization — months of work, orchestrated by agents.\n\nEli Lilly committed $1 billion over 5 years to NVIDIA-powered drug discovery. This is the same architecture."
    )


def slide_26_drug_detail(prs):
    _domain_detail_slide(prs,
        "ROADMAP — DRUG DISCOVERY", "Drug Discovery: BioNeMo Integration",
        [
            ("BioNeMo MolMIM", "Generative model for de novo molecular design with property constraints"),
            ("BioNeMo DiffDock", "Diffusion-based molecular docking — predict protein-ligand binding poses"),
            ("BioNeMo ESM-2", "Protein language model for target structure prediction and analysis"),
            ("ADMET Prediction", "Agent evaluates absorption, distribution, metabolism, excretion, toxicity"),
            ("Retrosynthesis Planning", "Agent plans synthetic routes for generated molecules"),
        ],
        "Lab DRUG3 builds a multi-agent pipeline: Hypothesis Agent generates candidates, Simulation Agent scores them,\n"
        "Analysis Agent filters by ADMET, Orchestrator Agent coordinates the full workflow.",
        notes="The detail: BioNeMo isn't just models — it's a complete platform. MolMIM generates molecules with desired properties. DiffDock predicts binding poses. ESM-2 analyzes protein targets.\n\nLab DRUG3 builds a full research crew: Hypothesis Agent generates candidates, Simulation Agent scores them, Analysis Agent filters by ADMET. One question in, a drug candidate pipeline out."
    )


def slide_27_materials(prs):
    _domain_overview_slide(prs,
        "ROADMAP — MATERIALS SCIENCE", "Materials Science: Accelerated Discovery with Physics-ML",
        "AI agents for property prediction, crystal design, and materials simulation",
        [
            ["MAT 1", "Property Prediction Agent", "Predict material properties using physics-ML surrogates"],
            ["MAT 2", "Crystal Structure Agent", "Generate and optimize crystal structures"],
            ["MAT 3", "Multi-Agent Materials Discovery", "Closed-loop: predict -> synthesize -> characterize"],
        ],
        ["Property\nPrediction", "Structure\nGeneration", "Simulation", "Characterize", "Optimize"],
        [CLR_MATERIAL, NVIDIA_GRAY, CLR_MATERIAL, NVIDIA_GRAY, CLR_MATERIAL],
        "NVIDIA SDKs: PhysicsNeMo (Modulus), RAPIDS (cuML, cuDF), NIM  |  Patterns: Simulation + Experiment Design",
        notes="Materials science is where PhysicsNeMo shines. Physics-informed neural networks embed conservation laws directly into the model training.\n\nThe result: simulations that run 1000x faster than finite element methods, with physical guarantees built in.\n\nThe agent invokes these surrogate models as tools — a researcher asks 'what's the thermal conductivity of this alloy?' and gets an answer in seconds instead of hours."
    )


def slide_28_materials_detail(prs):
    _domain_detail_slide(prs,
        "ROADMAP — MATERIALS SCIENCE", "Materials Science: PhysicsNeMo Integration",
        [
            ("Physics-Informed Neural Nets", "PINNs embed physical laws (conservation, boundary conditions) into training"),
            ("Neural Operators", "Fourier Neural Operators for surrogate models — 1000x faster than FEM"),
            ("RAPIDS cuML", "GPU-accelerated ML for materials property screening at scale"),
            ("Closed-Loop Discovery", "Agent generates candidates, simulates properties, ranks, and iterates"),
            ("Digital Twin Integration", "Connect to Omniverse for real-time materials simulation visualization"),
        ],
        "PhysicsNeMo provides pre-built architectures (FNO, PINN, DeepONet) that agents invoke as tools\n"
        "to run physics simulations in seconds instead of hours.",
        notes="PhysicsNeMo provides pre-built architectures — FNO, PINN, DeepONet — that agents invoke as tools. The key: simulations that respect physics laws but run 1000x faster.\n\nClosed-loop discovery means the agent generates candidates, simulates properties, ranks results, and iterates — all without human intervention on routine steps."
    )


def slide_29_climate(prs):
    _domain_overview_slide(prs,
        "ROADMAP — CLIMATE & EARTH SCIENCE", "Climate Science: AI Weather & Climate Agents",
        "Leverage Earth-2 models for forecasting, analysis, and climate research",
        [
            ["CLIM 1", "Weather Forecasting Agent", "Run AI forecasts with Earth-2 Medium Range model"],
            ["CLIM 2", "Climate Data Analysis Agent", "Analyze historical patterns, detect trends and anomalies"],
            ["CLIM 3", "Multi-Agent Climate Research", "Literature + forecast + analysis in coordinated pipeline"],
        ],
        ["Data\nIngestion", "AI Weather\nForecast", "Climate\nAnalysis", "Risk\nAssessment", "Report\nGeneration"],
        [NVIDIA_GRAY, CLR_CLIMATE, CLR_CLIMATE, NVIDIA_GRAY, CLR_CLIMATE],
        "NVIDIA SDKs: Earth-2 (Earth2Studio, CorrDiff, FourCastNet), RAPIDS, NIM  |  Patterns: Simulation + Data Analysis",
        notes="Earth-2 is NVIDIA's climate AI platform. The Medium Range model does 15-day weather forecasts in minutes. CorrDiff super-resolves from 25km to 2km.\n\nThe staggering number: AI can simulate 1000 years of climate in 24 hours. Traditional methods take weeks for 15-day forecasts.\n\nNOAA is already deploying AI weather models that use 0.3% of traditional compute."
    )


def slide_30_climate_detail(prs):
    _domain_detail_slide(prs,
        "ROADMAP — CLIMATE & EARTH SCIENCE", "Climate Science: Earth-2 Integration",
        [
            ("Earth-2 Medium Range", "15-day high-accuracy AI weather forecasts — minutes instead of hours"),
            ("Earth-2 Nowcasting", "Generative AI for 0-6 hour hazardous weather prediction"),
            ("CorrDiff Super-Resolution", "Downscale from 25km to 2km resolution using diffusion models"),
            ("Earth2Studio", "Python API for programmatic access to all Earth-2 models"),
            ("1000 Years in a Day", "AI can simulate 1000 years of climate in 24 hours vs. weeks"),
        ],
        "Agent uses Earth2Studio to run forecasts, RAPIDS cuDF for analysis at scale,\n"
        "and NIM for natural language interpretation of climate model outputs.",
        notes="The Earth-2 stack is comprehensive: Medium Range for 15-day forecasts, Nowcasting for 0-6 hour hazardous weather, CorrDiff for super-resolution.\n\nEarth2Studio provides the Python API that agents use as tools. Combined with RAPIDS for data processing and NIM for natural language interpretation, you get a complete climate research assistant."
    )


def slide_31_chemistry(prs):
    _domain_overview_slide(prs,
        "ROADMAP — CHEMISTRY", "Chemistry: Reaction Prediction & Process Optimization",
        "AI agents for synthesis planning, reaction prediction, and chemical engineering",
        [
            ["CHEM 1", "Reaction Prediction Agent", "Predict reaction outcomes and suggest conditions"],
            ["CHEM 2", "Process Optimization Agent", "Optimize parameters using cuOpt + PhysicsNeMo"],
            ["CHEM 3", "Multi-Agent Chemical Research", "Literature -> hypothesis -> simulate -> validate"],
        ],
        ["Literature\nReview", "Reaction\nPrediction", "Process\nSimulation", "Parameter\nOptimization", "Validation"],
        [NVIDIA_GRAY, CLR_CHEM, CLR_CHEM, NVIDIA_GRAY, CLR_CHEM],
        "NVIDIA SDKs: PhysicsNeMo, cuOpt (MILP/LP), RAPIDS, NIM  |  Patterns: Experiment Design + Simulation",
        notes="Chemistry combines PhysicsNeMo for molecular dynamics with cuOpt for process optimization.\n\nThe most compelling example: the Coscientist system at CMU designed and executed Nobel Prize-winning palladium-catalyzed reactions in under 4 minutes. That's the power of agent-orchestrated chemistry.\n\nOur lab builds toward that vision — literature search, hypothesis generation, simulation, optimization."
    )


def slide_32_chemistry_detail(prs):
    _domain_detail_slide(prs,
        "ROADMAP — CHEMISTRY", "Chemistry: PhysicsNeMo + cuOpt Integration",
        [
            ("PhysicsNeMo for MD", "Molecular dynamics simulation surrogates — accelerate by 100-1000x"),
            ("cuOpt for Scheduling", "GPU-accelerated optimization for reactor scheduling and process parameters"),
            ("Synthesis Planning", "Agent plans multi-step synthesis routes with retrosynthetic analysis"),
            ("Self-Driving Labs", "Agents that design, execute (via robotics), and analyze experiments 24/7"),
            ("Coscientist Pattern", "Nobel Prize-winning reactions designed and executed by AI in under 4 minutes"),
        ],
        "CHEM3 builds a research crew: Literature Agent finds prior art, Hypothesis Agent proposes reactions,\n"
        "Simulation Agent predicts outcomes, Optimization Agent tunes parameters.",
        notes="PhysicsNeMo provides molecular dynamics surrogates that accelerate simulations by 100-1000x. cuOpt handles GPU-accelerated optimization for reactor scheduling.\n\nThe self-driving lab vision is already real — agents design experiments, robots execute them, analysis agents interpret results, and the cycle repeats 24/7."
    )


def slide_33_robotics(prs):
    _domain_overview_slide(prs,
        "ROADMAP — ROBOTICS & LAB AUTOMATION", "Robotics: Self-Driving Laboratories",
        "AI agents that plan and execute physical experiments via robotic platforms",
        [
            ["ROB 1", "Lab Robot Planning Agent", "Translate protocols into executable robotic actions"],
            ["ROB 2", "Experiment Execution Agent", "Monitor, adapt, and recover during physical experiments"],
            ["ROB 3", "Multi-Agent Autonomous Lab", "Full loop: design -> execute -> analyze -> iterate"],
        ],
        ["Protocol\nDesign", "Robot\nPlanning", "Experiment\nExecution", "Result\nAnalysis", "Iterate"],
        [NVIDIA_GRAY, CLR_ROBOT, CLR_ROBOT, NVIDIA_GRAY, CLR_ROBOT],
        "NVIDIA SDKs: Isaac Sim, Omniverse, cuOpt (path planning), NIM  |  Patterns: Experiment Design + Multi-Agent",
        notes="Self-driving laboratories. Isaac Sim provides physically accurate robot simulation. cuOpt handles path planning. Omniverse creates digital twins of entire lab environments.\n\nThe pattern: agents design experiments, robots execute them, analysis agents interpret results, and the cycle repeats. Berkeley Lab's A-Lab is already doing this — AI proposes compounds, robots synthesize and test them.\n\nThe human stays in the loop for critical decisions."
    )


def slide_34_robotics_detail(prs):
    _domain_detail_slide(prs,
        "ROADMAP — ROBOTICS & LAB AUTOMATION", "Robotics: Isaac Sim + Omniverse Integration",
        [
            ("Isaac Sim", "Physically accurate robotic simulation — train and test before deploying"),
            ("Omniverse Digital Twins", "Create virtual replicas of lab environments for safe agent testing"),
            ("cuOpt Path Planning", "GPU-accelerated motion planning for robotic arms and mobile robots"),
            ("Protocol-to-Action", "Agent translates natural language lab protocols into robot commands"),
            ("Agent-to-Agent Workflows", "Multiple agents coordinate: planner, executor, analyzer, safety monitor"),
        ],
        "ROB3 demonstrates a full autonomous lab: agents design experiments, robots execute them,\n"
        "analysis agents interpret results, and the cycle repeats — all with human-in-the-loop approval gates.",
        notes="Isaac Sim lets you train and test robot behaviors in simulation before deploying to real hardware. Omniverse creates digital twins of entire lab environments.\n\nThe protocol-to-action translation is key — a scientist writes 'pipette 50uL of buffer A into well B3' and the agent converts that to robot commands. Multiple agents coordinate the full workflow with safety monitoring."
    )


def slide_35_finance(prs):
    _domain_overview_slide(prs,
        "DOMAIN — FINANCE (AVAILABLE)", "Finance: AI Agents for Quantitative Analysis & Risk",
        "Market analysis, portfolio risk assessment, and ESG compliance",
        [
            ["FIN 1", "Financial Analysis Agent", "Market data, VaR/CVaR risk, ESG screening + 4 visualizations"],
            ["FIN 2", "Portfolio Optimization Agent", "Mean-variance optimization, efficient frontier, rebalancing"],
            ["FIN 3", "Algo Strategy Agent", "Signal generation, backtesting, multi-agent trading pipeline"],
        ],
        ["Market\nData", "Performance\nAnalysis", "Risk\nAssessment", "ESG\nScreening", "Report\nGeneration"],
        [NVIDIA_GRAY, CLR_FINANCE, CLR_FINANCE, NVIDIA_GRAY, CLR_FINANCE],
        "NVIDIA SDKs: cuOpt (portfolio optimization), RAPIDS (GPU data science), NIM, Morpheus (fraud detection)",
        notes="Finance might surprise people in a science playbook, but quantitative finance IS computational science. Portfolio optimization, risk modeling, ESG analytics — all data-intensive, all benefiting from agents.\n\ncuOpt provides GPU-accelerated optimization. RAPIDS processes millions of trades per second. Morpheus detects fraud in real time.\n\nAnd NIM's low latency — sub-100ms — is critical for time-sensitive trading decisions."
    )


def slide_36_finance_detail(prs):
    _domain_detail_slide(prs,
        "DOMAIN — FINANCE (AVAILABLE)", "Finance: NVIDIA cuOpt + RAPIDS Integration",
        [
            ("cuOpt for Portfolio Optimization", "GPU-accelerated MILP/LP solvers for optimal asset allocation"),
            ("RAPIDS cuDF", "Process millions of trades per second — GPU-accelerated pandas"),
            ("Morpheus for Fraud Detection", "Real-time anomaly detection in transaction streams"),
            ("Low-Latency NIM", "Sub-100ms inference for time-critical trading decisions"),
            ("ESG Analytics", "Agent screens portfolios against sustainability frameworks (SASB, TCFD, GRI)"),
        ],
        "FIN1 builds a complete quant analysis agent: market data -> risk assessment -> ESG screening.\n"
        "The same agent pattern scales from research notebooks to production trading systems.",
        notes="cuOpt provides GPU-accelerated MILP/LP solvers — the same optimization engines used in production trading systems, but accessible through a simple agent tool call.\n\nRAPIDS cuDF is GPU-accelerated pandas — process millions of rows in milliseconds. Combined with Morpheus for real-time fraud detection, you get a complete fintech agent stack."
    )


# ── SECTION 7: NVIDIA SDK DEEP DIVE (slides 35-38) ────────────────────────

def slide_35_nim_deep(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_common(slide, "NVIDIA SDK — NIM")
    _add_title(slide, "Why NVIDIA NIM for Scientific Agents?")
    _add_subtitle(slide, "Not just HOW to switch — here's WHY you'd choose NIM over OpenAI")
    # Decision table - the WHY
    _add_table(slide,
        ["Factor", "OpenAI API", "NVIDIA NIM"],
        [
            ["Data Privacy", "Data sent to OpenAI servers", "Data stays in YOUR infrastructure"],
            ["Cost at Scale", "Pay per token (expensive at 10K+ calls)", "Pay per GPU hour (self-hosted) — 3-10x cheaper"],
            ["Latency", "~500ms typical", "~200ms on GPU (optimized inference)"],
            ["Models", "GPT-4o (closed)", "Nemotron, Llama 3.3, 100+ open models"],
            ["Fine-tuning", "Limited", "Full control with NeMo Framework"],
            ["Compliance", "Cloud only", "On-prem for HIPAA, pharma, finance regulations"],
        ],
        top=Inches(1.6), row_h=Pt(30)
    )
    _add_code_block(slide,
        '# Switch in one line — zero code changes:\n'
        'export USE_NIM=true\n'
        'export NIM_API_KEY="nvapi-..."   # Free at build.nvidia.com',
        top=Inches(5.5), height=Inches(1.0)
    )
    _add_box(slide, "build.nvidia.com — Free tier  |  100+ models  |  API key in 30 seconds  |  On-prem deployment available",
             Inches(0.6), Inches(6.7), Inches(12), Inches(0.4),
             fill_color=NVIDIA_GREEN, text_color=NVIDIA_BLACK, font_size=Pt(13), bold=True)
    _add_notes(slide, "This is the slide that matters for your infrastructure team. Left column: OpenAI. Right column: NVIDIA NIM.\n\nThree things jump out: data privacy — your patient data, proprietary molecules, unpublished results stay in YOUR infrastructure. Cost at scale — self-hosted NIM is 3-10x cheaper when you're running thousands of agent calls. And compliance — on-prem deployment for HIPAA, pharma regulations, financial compliance.\n\nThe switch? One environment variable. Zero code changes.\n\n[PAUSE] What are your data privacy requirements? That usually drives the NIM conversation.")


def slide_36_domain_sdks(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_common(slide, "NVIDIA SDK — DOMAIN ACCELERATION")
    _add_title(slide, "Domain-Specific NVIDIA SDKs")
    _add_table(slide,
        ["SDK", "Domain", "Key Capabilities", "Used In"],
        [
            ["BioNeMo", "Drug Discovery, Genomics", "MolMIM, DiffDock, ESM-2, protein folding", "DRUG, BIO labs"],
            ["Earth-2", "Climate & Weather", "Medium Range, Nowcasting, CorrDiff, Earth2Studio", "CLIM labs"],
            ["PhysicsNeMo", "Materials, Chemistry, Physics", "PINNs, FNO, DeepONet, neural operators", "MAT, CHEM labs"],
            ["RAPIDS", "All domains", "cuDF, cuML, cuGraph — GPU data science", "All domain labs"],
            ["cuOpt", "Energy, Chemistry, Robotics", "GPU-accelerated MILP/LP/VRP optimization", "CHEM, ROB labs"],
            ["Isaac Sim", "Robotics", "Physically accurate robot simulation", "ROB labs"],
            ["Omniverse", "Physics, Robotics", "Digital twins, real-time visualization", "ROB, MAT labs"],
            ["cuOpt + RAPIDS", "Finance", "Portfolio optimization, GPU data processing, fraud detection", "FIN labs"],
        ],
        top=Inches(1.6), row_h=Pt(32)
    )
    _add_notes(slide, "Quick reference for your technical team. Seven NVIDIA SDKs mapped to the domains where they apply.\n\nBioNeMo for drug discovery and genomics. Earth-2 for climate. PhysicsNeMo for materials and chemistry. RAPIDS everywhere — it's GPU-accelerated pandas.\n\nWhich of these domains is most relevant to your organization? [PAUSE for customer input]")


def slide_37_guardrails(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_common(slide, "NVIDIA SDK — SAFETY")
    _add_title(slide, "NeMo Guardrails: Safe Scientific Agents")
    _add_subtitle(slide, "Every scientific domain has unique safety requirements")
    items = [
        ("Fabricated Citations", "Guardrails verify that cited papers actually exist in the literature"),
        ("Chemical Safety", "Rails prevent agents from suggesting dangerous synthesis routes"),
        ("Patient Data Privacy", "PII detection and masking for clinical notes and genomic data"),
        ("Experiment Safety", "Approval gates before agents trigger expensive simulations or physical experiments"),
        ("Hallucination Detection", "Fact-checking rails ground claims against retrieved evidence"),
        ("Reproducibility Enforcement", "EOP rails ensure agents document evidence chains for every claim"),
    ]
    _add_bullets(slide, items, top=Inches(1.7), size=Pt(15))
    _add_notes(slide, "Safety isn't optional for scientific agents. Every domain has unique risks.\n\nIn healthcare: patient data privacy and drug interaction checking. In chemistry: preventing dangerous synthesis routes. In genomics: ensuring variant classifications are evidence-based, not hallucinated.\n\nNeMo Guardrails is open-source and integrates directly with the agent loop. It's not bolted on — it's baked in.")


def slide_38_agent_toolkit(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_common(slide, "NVIDIA SDK — ORCHESTRATION")
    _add_title(slide, "NeMo Agent Toolkit: Framework-Agnostic Orchestration")
    _add_subtitle(slide, "Build once, run on any framework — config-driven agent workflows")
    _add_code_block(slide,
        '# config.yml — declarative agent workflow\n'
        'workflow:\n'
        '  name: scientific_research_agent\n'
        '  llm:\n'
        '    provider: nim\n'
        '    model: nvidia/llama-3.3-nemotron-super-49b-v1.5\n'
        '  tools:\n'
        '    - pubmed_search\n'
        '    - molecular_generator\n'
        '    - docking_scorer\n'
        '  guardrails:\n'
        '    - citation_verification\n'
        '    - chemical_safety\n'
        '  framework: langgraph   # or: langchain, llamaindex, crewai\n'
        '  evaluation:\n'
        '    judge_model: nvidia/nemotron-70b\n'
        '    rubric: scientific_accuracy',
        top=Inches(1.6), height=Inches(4.8)
    )
    _add_box(slide, "Supports: LangChain  |  LlamaIndex  |  CrewAI  |  Semantic Kernel  |  Google ADK",
             Inches(0.6), Inches(6.6), Inches(12), Inches(0.4),
             fill_color=NVIDIA_DARK, text_color=NVIDIA_GREEN, font_size=Pt(13), bold=True)
    _add_notes(slide, "The NeMo Agent Toolkit is NVIDIA's framework for production agents. The key advantage: it's config-driven and framework-agnostic.\n\nYour team defines workflows in YAML — which LLM, which tools, which guardrails, which evaluation rubric. Then they can swap from LangChain to LlamaIndex to CrewAI without changing code.\n\nThis is how you go from notebook prototype to production deployment.")


# ── SECTION 8: AGENT PATTERNS (slides 39-41) ──────────────────────────────

def slide_39_patterns_catalog(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_common(slide, "AGENT PATTERNS")
    _add_title(slide, "Agent Pattern Catalog for Scientific Discovery")
    _add_table(slide,
        ["Pattern", "Description", "Best For", "Example Domain"],
        [
            ["Hypothesis Generation", "Creative reasoning + literature grounding", "Generating novel research ideas", "Drug Discovery"],
            ["Literature Review (RAG)", "Search, retrieve, evaluate, synthesize", "Evidence-based summaries", "Healthcare"],
            ["Experiment Design", "Optimize parameters, generate protocols", "Planning lab experiments", "Chemistry"],
            ["Data Analysis", "Statistical analysis + anomaly detection", "Making sense of large datasets", "Climate Science"],
            ["Simulation", "Invoke surrogate models, interpret results", "Physics-ML prediction", "Materials Science"],
            ["Multi-Agent Orchestrator", "Coordinate specialized agent teams", "End-to-end workflows", "All domains"],
        ],
        top=Inches(1.6)
    )
    _add_notes(slide, "Six patterns, each one reusable across domains. The Hypothesis Generation pattern works for drug discovery AND materials science. The RAG pattern works for healthcare literature AND climate research.\n\nYour team learns these once in the Foundation, then applies them everywhere.")


def slide_40_pattern_mapping(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_common(slide, "AGENT PATTERNS")
    _add_title(slide, "Pattern-to-Domain Mapping")
    _add_subtitle(slide, "Primary = dedicated notebook  |  Secondary = supporting role")
    _add_table(slide,
        ["Pattern", "EOP", "Health", "Bio", "Drug", "Mat", "Climate", "Chem", "Robot", "Fin"],
        [
            ["Hypothesis Gen", "-", "-", "P", "P", "S", "S", "S", "-", "S"],
            ["Literature RAG", "P", "P", "P", "S", "P", "P", "P", "S", "S"],
            ["Experiment Design", "-", "-", "S", "S", "P", "-", "P", "P", "-"],
            ["Data Analysis", "P", "P", "P", "S", "P", "P", "P", "S", "P"],
            ["Simulation", "-", "-", "-", "P", "P", "P", "P", "P", "S"],
            ["Multi-Agent", "S", "S", "P", "P", "S", "S", "S", "P", "S"],
        ],
        top=Inches(1.8)
    )
    _add_box(slide, "P = Primary (dedicated notebook in that lab)     S = Secondary (supporting role)     - = Not featured",
             Inches(0.6), Inches(5.8), Inches(12), Inches(0.4),
             fill_color=NVIDIA_DARK, text_color=NVIDIA_LGRAY, font_size=Pt(12))
    _add_notes(slide, "This matrix shows where each pattern is used as a primary feature vs. supporting role. The takeaway: every domain uses at least 3-4 patterns. Your team builds a vocabulary of agent patterns that compounds over time.")


def slide_41_multi_agent(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_common(slide, "AGENT PATTERNS")
    _add_title(slide, "Multi-Agent Orchestration: The Scientific Research Crew")
    _add_subtitle(slide, "Specialized agents collaborate on complex research workflows")
    agents = [
        ("ORCHESTRATOR", "Decomposes goals,\nroutes tasks,\naggregates results", NVIDIA_GREEN),
        ("Literature\nAgent", "Searches papers,\nevaluates evidence,\nsummarizes findings", CLR_HEALTH),
        ("Hypothesis\nAgent", "Generates novel\nideas, scores\nnovelty", CLR_DRUG),
        ("Simulation\nAgent", "Runs physics-ML\nmodels, interprets\nresults", CLR_MATERIAL),
        ("Analysis\nAgent", "Statistical tests,\nvisualization,\nanomaly detection", CLR_CLIMATE),
    ]
    # Orchestrator in center top
    _add_box(slide, f"{agents[0][0]}\n{agents[0][1]}", Inches(4.8), Inches(1.8), Inches(3.5), Inches(1.6),
             fill_color=agents[0][2], text_color=NVIDIA_BLACK, font_size=Pt(12), bold=True)
    # 4 specialist agents below
    for i, (name, desc, color) in enumerate(agents[1:]):
        left = Inches(0.3) + Inches(i * 3.25)
        _add_box(slide, f"{name}\n{desc}", left, Inches(4.5), Inches(3.0), Inches(1.8),
                 fill_color=color, text_color=NVIDIA_WHITE, font_size=Pt(11), bold=True)
        # Arrow from orchestrator
        _add_arrow_down(slide, Inches(1.5) + Inches(i * 3.25), Inches(3.5), Inches(0.25), Inches(0.9))
    _add_notes(slide, "The multi-agent orchestrator is the most powerful pattern. Specialized agents — each focused on one task — coordinated by an orchestrator.\n\nThis mirrors how real research teams work. A literature specialist, a hypothesis generator, a simulation expert, a data analyst. Each agent is excellent at one thing. The orchestrator makes them work together.\n\nThis is what Lab 6 teaches — and what production systems like NVIDIA's drug discovery pipelines use at scale.")


# ── SECTION 9: PRODUCTION & CLOSING (slides 42-50) ────────────────────────

def slide_42_production_path(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_common(slide, "PRODUCTION PATH")
    _add_title(slide, "From Notebook to Production")
    steps = [
        ("1. PROTOTYPE", "Jupyter notebook\nNIM API\nSingle agent", NVIDIA_GRAY),
        ("2. EVALUATE", "LLM-as-Judge\nBenchmark suite\nA/B testing", NVIDIA_GRAY),
        ("3. HARDEN", "NeMo Guardrails\nError recovery\nHuman-in-loop", NVIDIA_GREEN2),
        ("4. DEPLOY", "Triton Inference\nKubernetes\nHelm charts", NVIDIA_GREEN),
        ("5. MONITOR", "Agent observability\nLatency tracking\nAudit logging", NVIDIA_GREEN),
    ]
    for i, (title, desc, color) in enumerate(steps):
        left = Inches(0.3) + Inches(i * 2.6)
        _add_box(slide, f"{title}\n\n{desc}", left, Inches(1.8), Inches(2.3), Inches(3.0),
                 fill_color=color, text_color=NVIDIA_WHITE, font_size=Pt(12), bold=True)
        if i < 4:
            _add_arrow_right(slide, left + Inches(2.35), Inches(3.0), Inches(0.2), Inches(0.2))
    _add_box(slide, "Part 5 of the playbook covers Triton serving, Kubernetes deployment, monitoring, and three end-to-end case studies",
             Inches(0.3), Inches(5.3), Inches(12.6), Inches(0.5),
             fill_color=NVIDIA_DARK, text_color=NVIDIA_LGRAY, font_size=Pt(13))
    _add_notes(slide, "Five stages: prototype in a notebook, evaluate with LLM-as-Judge, harden with guardrails and error recovery, deploy on Triton and Kubernetes, monitor with observability.\n\nMost teams get to stage 2 in a week. Stage 3 takes another week. Stages 4-5 depend on your infrastructure — but NIM makes it straightforward.")


def slide_43_human_in_loop(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_common(slide, "DESIGN PRINCIPLES")
    _add_title(slide, "Human-in-the-Loop: Scientists Stay in Control")
    items = [
        ("Approval Gates", "Before expensive simulations ($), dangerous experiments, or irreversible actions"),
        ("Confidence Thresholds", "Agent auto-proceeds when confidence > 90%, asks human when uncertain"),
        ("Interactive Review", "Scientist can modify agent-generated hypotheses before testing"),
        ("Audit Trails", "Every agent decision logged: what, when, why, which tool, what result"),
        ("Escalation Paths", "Agent detects when it's out of depth and requests human expertise"),
    ]
    _add_bullets(slide, items, top=Inches(1.7), size=Pt(16))
    _add_box(slide, "Key Principle: Agents automate the routine. Scientists make the critical decisions.",
             Inches(0.6), Inches(5.8), Inches(12), Inches(0.5),
             fill_color=NVIDIA_GREEN, text_color=NVIDIA_BLACK, font_size=Pt(15), bold=True)
    _add_notes(slide, "This matters for regulated industries. Scientists stay in control.\n\nApproval gates before expensive operations. Confidence thresholds — the agent auto-proceeds when confident, asks a human when uncertain. Full audit trails.\n\nThe key principle: agents automate the routine. Scientists make the critical decisions.")


def slide_44_extensibility(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_common(slide, "EXTENSIBILITY")
    _add_title(slide, "Add Your Own Domain in 3 Labs")
    _add_code_block(slide,
        '# Contributing a new domain is simple:\n'
        '#\n'
        '# domains/\n'
        '#   your_domain/\n'
        '#     Lab_YD1_Introduction.ipynb        # Domain SDK + basic agent\n'
        '#     Lab_YD2_Core_Workflow.ipynb        # Main pipeline with tools\n'
        '#     Lab_YD3_Multi_Agent.ipynb          # Orchestrated multi-agent\n'
        '#\n'
        '# Each lab follows the template:\n'
        '#   1. Setup (pip install + make_client)\n'
        '#   2. Background (domain context for non-experts)\n'
        '#   3. Tool Schemas (Pydantic definitions)\n'
        '#   4. Tool Implementations (Python + LLM-powered)\n'
        '#   5. Agent Function (system prompt + tool calling)\n'
        '#   6. Experiments (3-5 demonstrations)\n'
        '#   7. Summary (tools built, patterns used, next steps)',
        top=Inches(1.5), height=Inches(4.5)
    )
    domains = ["Astronomy", "Energy", "Ecology", "Neuroscience", "Your\nDomain"]
    for i, d in enumerate(domains):
        left = Inches(0.5) + Inches(i * 2.55)
        color = NVIDIA_GREEN if i == 4 else NVIDIA_GRAY
        _add_box(slide, d, left, Inches(6.3), Inches(2.2), Inches(0.7),
                 fill_color=color, text_color=NVIDIA_WHITE, font_size=Pt(13), bold=True)
    _add_notes(slide, "Your team isn't limited to our domains. The 3-lab template means any team can add their scientific domain in a few days.\n\nAstronomy, ecology, neuroscience — whatever your organization needs. The foundation patterns work everywhere.")


def slide_45_beginner_friendly(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_common(slide, "DESIGNED FOR BEGINNERS")
    _add_title(slide, "Beginner-Friendly Features Built Into Every Lab")
    _add_subtitle(slide, "We tested with beginners and fixed the top 4 barriers to learning")
    improvements = [
        ("Gentle LangGraph On-Ramp",
         "Lab 4 now starts with a 2-node warm-up graph before the full research pipeline. "
         "No more \"cliff\" from Python classes to state machines."),
        ("Genomics Primer for Non-Biologists",
         "BIO Lab 2 includes a 5-minute crash course: DNA→Protein, variant types, ACMG basics, "
         "and how to read variant notation — so anyone can follow along."),
        ("NIM Decision Guide in Lab 0",
         "Clear comparison table: WHEN to use OpenAI (learning, prototyping) vs WHEN to use NIM "
         "(data privacy, cost control, regulated industries, open models)."),
        ("Expected Output for Every Experiment",
         "61 collapsible sample outputs across all 18 notebooks — read-only learners can follow "
         "the full playbook without running any code."),
        ("Step-by-Step Code Explanations",
         "Every code cell has a preceding markdown cell explaining WHAT it does in 2-4 sentences. "
         "Large code blocks split into small, digestible pieces."),
    ]
    for i, (title, desc) in enumerate(improvements):
        tp = Inches(1.7) + Inches(i * 1.0)
        _add_box(slide, str(i + 1), Inches(0.4), tp, Inches(0.45), Inches(0.45),
                 fill_color=NVIDIA_GREEN, text_color=NVIDIA_BLACK, font_size=Pt(18), bold=True)
        txBox = slide.shapes.add_textbox(Inches(1.1), tp - Inches(0.02), Inches(11.5), Inches(0.3))
        p = txBox.text_frame.paragraphs[0]
        p.text = title; p.font.size = Pt(16); p.font.color.rgb = NVIDIA_GREEN; p.font.bold = True; p.font.name = FONT_BODY
        txBox2 = slide.shapes.add_textbox(Inches(1.1), tp + Inches(0.32), Inches(11.5), Inches(0.5))
        tf2 = txBox2.text_frame; tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = desc; p2.font.size = Pt(12); p2.font.color.rgb = NVIDIA_LGRAY; p2.font.name = FONT_BODY
    _add_notes(slide, "We didn't just build labs — we tested them with beginners and fixed the top barriers.\n\nLangGraph warm-up eliminates the learning cliff. Genomics primer makes BIO labs accessible to non-biologists. NIM decision guide answers 'why would I switch from OpenAI?' Expected output means you can read the whole playbook without running code.\n\nYour team doesn't need to be AI experts to start. They need curiosity and Python.")


def slide_46_case_studies(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_common(slide, "CASE STUDIES")
    _add_title(slide, "Real-World Impact: Agentic AI in Science")
    items = [
        ("Coscientist (CMU)", "AI designed and executed Nobel Prize-winning palladium-catalyzed reactions in < 4 minutes"),
        ("A-Lab (Berkeley)", "AI proposes new compounds, robots prepare and test them in closed-loop discovery"),
        ("NOAA AI Weather", "AI-driven global forecast uses 0.3% of traditional compute, 16-day prediction in 40 min"),
        ("Eli Lilly + BioNeMo", "$1B commitment — AI agents accelerating drug discovery pipeline"),
        ("Google AI Co-Scientist", "Multi-agent system generates and validates novel research hypotheses"),
    ]
    _add_bullets(slide, items, top=Inches(1.7), size=Pt(16))
    _add_notes(slide, "This isn't theoretical. Coscientist at CMU executed Nobel Prize-winning reactions autonomously. Berkeley Lab's A-Lab runs closed-loop materials discovery. NOAA's AI weather models cut compute by 99.7%. Eli Lilly committed a billion dollars.\n\nThe question isn't whether AI agents work for science. It's how quickly your organization starts building them.")


def slide_46_demo(prs):
    _section_divider(prs, "LIVE DEMO",
                     "Lab 0: Build an Agent from Scratch\n"
                     "Lab HC2: Medical Literature Agent\n"
                     "Lab DRUG1: Molecular Generation with BioNeMo",
                     notes="[If doing live demo] Let me show you this live. I'll run Lab 0 — build an agent in 3 minutes.\n\n[If not doing live demo] In the interest of time, I'll skip the live demo — but your team will run this in their first 10 minutes with the playbook.")


def slide_47_getting_started(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_common(slide, "GETTING STARTED")
    _add_title(slide, "Get Started in 5 Minutes")
    _add_code_block(slide,
        '# 1. Clone the repository\n'
        'git clone https://github.com/NVIDIA/agentic-ai-science-playbook.git\n'
        'cd agentic-ai-science-playbook\n'
        '\n'
        '# 2. Install dependencies\n'
        'pip install -r requirements.txt\n'
        '\n'
        '# 3. Set your API key (choose one)\n'
        'export OPENAI_API_KEY="sk-..."                # OpenAI\n'
        '# OR\n'
        'export USE_NIM=true                            # NVIDIA NIM\n'
        'export NIM_API_KEY="nvapi-..."                 # Free at build.nvidia.com\n'
        '\n'
        '# 4. Start with the Scenarios Lab (no coding needed!)\n'
        'jupyter lab Lab_Scenarios_AI_Agent_Research.ipynb\n'
        '\n'
        '# 5. Then build your own agents\n'
        'jupyter lab Lab0_Agent_Prototype.ipynb         # Foundation\n'
        'jupyter lab Lab6_AI_CoScientist_MultiAgent.ipynb  # Multi-agent',
        top=Inches(1.5), height=Inches(5.0)
    )
    _add_box(slide, "No GPU required  |  Google Colab ready  |  Python 3.10+  |  NVIDIA NIM free tier",
             Inches(0.6), Inches(6.7), Inches(12), Inches(0.4),
             fill_color=NVIDIA_GREEN, text_color=NVIDIA_BLACK, font_size=Pt(13), bold=True)
    _add_notes(slide, "Five commands. That's it. Clone, install, set your API key, open the Scenarios Lab.\n\nNo GPU required. Works on any laptop. Free NVIDIA NIM API key at build.nvidia.com.\n\nYour team can start this afternoon.\n\n[PAUSE] Want us to set up a workshop for your team? We can do a half-day guided session through the Foundation Labs.")


def slide_48_takeaways(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_common(slide, "KEY TAKEAWAYS")
    _add_title(slide, "Key Takeaways")
    items = [
        ("Agent engineering is a learnable skill, not magic",
         "7 foundation labs + Scenarios tutorial take you from zero to AI Co-Scientist"),
        ("Beginner-friendly by design",
         "Step-by-step explanations, expected output, LangGraph warm-up, genomics primer, NIM decision guide"),
        ("NVIDIA NIM: data privacy, speed, and open models",
         "Keep sensitive research data on-prem, 2-5x faster inference, full fine-tuning control"),
        ("Foundation patterns transfer across ALL 5+ scientific domains",
         "Same agent loop powers healthcare, genomics, finance, climate, drug discovery, and more"),
    ]
    for i, (title, desc) in enumerate(items):
        tp = Inches(1.6) + Inches(i * 1.4)
        _add_box(slide, str(i + 1), Inches(0.5), tp, Inches(0.5), Inches(0.5),
                 fill_color=NVIDIA_GREEN, text_color=NVIDIA_BLACK, font_size=Pt(20), bold=True)
        txBox = slide.shapes.add_textbox(Inches(1.3), tp - Inches(0.02), Inches(11), Inches(0.4))
        p = txBox.text_frame.paragraphs[0]
        p.text = title; p.font.size = Pt(18); p.font.color.rgb = NVIDIA_WHITE; p.font.bold = True; p.font.name = FONT_BODY
        txBox2 = slide.shapes.add_textbox(Inches(1.3), tp + Inches(0.42), Inches(11), Inches(0.4))
        p2 = txBox2.text_frame.paragraphs[0]
        p2.text = desc; p2.font.size = Pt(14); p2.font.color.rgb = NVIDIA_LGRAY; p2.font.name = FONT_BODY
    _add_notes(slide, "Four things to remember.\n\nFirst: agent engineering is learnable. Your team doesn't need PhDs in AI — they need Python and curiosity.\n\nSecond: we built this for beginners. No cliff, no jargon walls.\n\nThird: NVIDIA NIM gives you data privacy, speed, and open models. Your research stays yours.\n\nFourth: these patterns work everywhere. Learn once, apply across every domain in your organization.")


def slide_49_resources(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_common(slide, "RESOURCES")
    _add_title(slide, "Resources & Next Steps")
    items = [
        ("Playbook Repository", "github.com/NVIDIA/agentic-ai-science-playbook"),
        ("NVIDIA NIM", "build.nvidia.com — free API access, 100+ models"),
        ("NVIDIA BioNeMo", "nvidia.com/bionemo — drug discovery & genomics models"),
        ("NVIDIA Earth-2", "nvidia.com/earth-2 — climate & weather AI"),
        ("NVIDIA PhysicsNeMo", "developer.nvidia.com/physicsnemo — physics-ML models"),
        ("NeMo Agent Toolkit", "developer.nvidia.com/nemo-agent-toolkit — agent framework"),
        ("NeMo Guardrails", "developer.nvidia.com/nemo-guardrails — safety toolkit"),
        ("NVIDIA DLI Courses", "nvidia.com/training — agent AI certifications"),
    ]
    _add_bullets(slide, items, top=Inches(1.6), size=Pt(15))
    _add_notes(slide, "All the links you need. The playbook repo, NVIDIA NIM, the domain SDKs, DLI courses for certification.\n\nI'll share these slides with the deck — all links are clickable.")


def slide_50_thankyou(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(0.15))
    bar.fill.solid(); bar.fill.fore_color.rgb = NVIDIA_GREEN; bar.line.fill.background()
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(3), Inches(0.5))
    p = txBox.text_frame.paragraphs[0]
    p.text = "NVIDIA"; p.font.size = Pt(18); p.font.color.rgb = NVIDIA_GREEN; p.font.bold = True; p.font.name = FONT_TITLE
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.3), Inches(11), Inches(1.5))
    p = txBox.text_frame.paragraphs[0]
    p.text = "Thank You"; p.font.size = Pt(52); p.font.color.rgb = NVIDIA_WHITE; p.font.bold = True; p.font.name = FONT_TITLE; p.alignment = PP_ALIGN.CENTER
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(11), Inches(0.8))
    p2 = txBox2.text_frame.paragraphs[0]
    p2.text = "Questions?"; p2.font.size = Pt(28); p2.font.color.rgb = NVIDIA_GREEN; p2.font.name = FONT_BODY; p2.alignment = PP_ALIGN.CENTER
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(4.8), Inches(4), Inches(0.04))
    line.fill.solid(); line.fill.fore_color.rgb = NVIDIA_GREEN; line.line.fill.background()
    txBox3 = slide.shapes.add_textbox(Inches(1), Inches(5.1), Inches(11), Inches(1.2))
    tf3 = txBox3.text_frame; tf3.word_wrap = True
    for i, txt in enumerate([
        "github.com/NVIDIA/agentic-ai-science-playbook",
        "build.nvidia.com  |  developer.nvidia.com",
        "NVIDIA NIM  |  BioNeMo  |  Earth-2  |  PhysicsNeMo  |  NeMo Agent Toolkit",
    ]):
        p = tf3.paragraphs[0] if i == 0 else tf3.add_paragraph()
        p.text = txt; p.font.size = Pt(14); p.font.color.rgb = NVIDIA_LGRAY; p.font.name = FONT_BODY
        p.alignment = PP_ALIGN.CENTER; p.space_after = Pt(4)
    _add_bottom_bar(slide)
    _add_notes(slide, "Thank you for your time. I'm genuinely excited about what your team can build with this.\n\nQuestions? I'm happy to dive deeper into any domain or discuss how to run a workshop for your team.\n\n[Have backup slides ready for: pricing, on-prem deployment, comparison with other agent frameworks, compliance/regulatory questions]")


# ═══════════════════════════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════════════════════════

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Section 1: Title & Intro (4 slides) — establish problem BEFORE scenarios
    slide_01_title(prs)
    slide_02_agenda(prs)
    slide_03_roadmap(prs)
    slide_05_landscape(prs)
    slide_04_scenarios(prs)

    # Section 2: Agent Architecture (5 slides)
    _section_divider(prs, "LLM Agent Architecture", "Understanding the building blocks of scientific AI agents",
                     notes="Now let's get under the hood. What makes an agent fundamentally different from ChatGPT?")
    slide_04_what_is_agent(prs)
    slide_05_agent_components(prs)
    slide_06_why_science(prs)
    slide_07_chatbot_vs_agent(prs)
    slide_08_nvidia_stack(prs)

    # Section 3: Playbook Architecture (3 slides)
    _section_divider(prs, "Playbook Architecture", "Foundation + 9 scientific domains",
                     notes="Now let's look at how the playbook is organized — foundation plus domains.")
    slide_09_architecture(prs)
    slide_10_learning_paths(prs)
    slide_11_tech_stack(prs)

    # Section 4: Foundation Labs (8 slides)
    _section_divider(prs, "Foundation Labs", "7 production patterns every scientific agent needs",
                     notes="This is the core of the playbook. 7 labs, each building on the last, taking your team from zero to production.")
    slide_12_foundation_overview(prs)
    slide_13_lab0(prs)
    slide_14_lab1_lab2(prs)
    slide_15_lab3(prs)
    slide_16_lab4(prs)
    slide_17_lab5(prs)
    slide_18_foundation_summary(prs)
    slide_19_coscientist(prs)

    # Section 5: Available Domain Labs (8 slides — EOP, Healthcare, Bio, Finance)
    _section_divider(prs, "Domain Labs: Available Today", "EOP (3 Labs)  |  Healthcare (3 Labs)  |  Bioinformatics (3 Labs)  |  Finance (3 Labs)",
                     notes="These four domain tracks are fully built and tested. 10 domain notebooks your team can run today, plus the 8 foundation notebooks — 18 total.")
    slide_20_eop(prs)
    slide_21_eop_detail(prs)
    slide_22_healthcare(prs)
    slide_23_healthcare_detail(prs)
    slide_24_bioinformatics(prs)
    slide_25_bioinformatics_detail(prs)
    slide_35_finance(prs)
    slide_36_finance_detail(prs)

    # Section 6: Roadmap Domain Labs (10 slides — planned NVIDIA SDK integrations)
    _section_divider(prs, "Domain Roadmap", "Planned NVIDIA SDK Integrations:  Drug Discovery  |  Materials  |  Climate  |  Chemistry  |  Robotics",
                     notes="Now let me show you what's coming. These 5 domains are on our roadmap — each one integrates a specific NVIDIA domain SDK. The foundation patterns your team already learned will apply directly. Let me highlight the ones most relevant to your organization.")
    slide_25_drug_discovery(prs)
    slide_26_drug_detail(prs)
    slide_27_materials(prs)
    slide_28_materials_detail(prs)
    slide_29_climate(prs)
    slide_30_climate_detail(prs)
    slide_31_chemistry(prs)
    slide_32_chemistry_detail(prs)
    slide_33_robotics(prs)
    slide_34_robotics_detail(prs)

    # Section 7: NVIDIA SDK Deep Dive (4 slides)
    _section_divider(prs, "NVIDIA SDK Deep Dive", "NIM  |  Domain SDKs  |  Guardrails  |  Agent Toolkit",
                     notes="Let's go deeper on NVIDIA's SDKs — specifically, WHY you'd choose NIM for your research infrastructure.")
    slide_35_nim_deep(prs)
    slide_36_domain_sdks(prs)
    slide_37_guardrails(prs)
    slide_38_agent_toolkit(prs)

    # Section 8: Agent Patterns (3 slides)
    _section_divider(prs, "Agent Pattern Catalog", "6 reusable patterns mapped across 9 domains",
                     notes="Let me show you the 6 reusable patterns — and how they map across all domains.")
    slide_39_patterns_catalog(prs)
    slide_40_pattern_mapping(prs)
    slide_41_multi_agent(prs)

    # Section 9: Production & Closing (10 slides)
    _section_divider(prs, "Production & Getting Started", "From notebook to deployment",
                     notes="Finally — how do you go from learning to production?")
    slide_42_production_path(prs)
    slide_43_human_in_loop(prs)
    slide_44_extensibility(prs)
    slide_45_beginner_friendly(prs)
    slide_46_case_studies(prs)
    slide_46_demo(prs)
    slide_47_getting_started(prs)
    slide_48_takeaways(prs)
    slide_49_resources(prs)
    slide_50_thankyou(prs)

    out_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "Agentic_AI_Science_Playbook.pptx")
    prs.save(out_path)
    print(f"Saved {len(prs.slides)} slides -> {out_path}")


if __name__ == "__main__":
    main()
